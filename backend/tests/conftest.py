import pytest


@pytest.fixture
def data_dir(tmp_path):
    return tmp_path / "data"


@pytest.fixture
def sample_pptx_path(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Test Slide"

    path = tmp_path / "test.pptx"
    prs.save(str(path))

    return str(path)


@pytest.fixture
def sample_excel_path(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    ws.append(["Category", "Value", "Growth"])
    ws.append(["Revenue", 1000, 0.25])
    ws.append(["Costs", 750, -0.10])
    ws.append(["Profit", 250, 0.35])

    path = tmp_path / "test.xlsx"
    wb.save(str(path))

    return str(path)
