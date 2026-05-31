import os
from typing import Optional
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook


from .model_registry import model_registry

_logger_doc = None
def _get_doc_logger():
    global _logger_doc
    if _logger_doc is None:
        import logging
        _logger_doc = logging.getLogger("slideforge.document_ingestion")
    return _logger_doc


@dataclass
class TextRun:
    text: str
    font_name: str
    font_size: float
    font_bold: bool
    font_italic: bool
    font_color_rgb: tuple[int, int, int]
    shape_id: str
    text_box_id: str
    start_index: int
    end_index: int


@dataclass
class TextBox:
    id: str
    text: str
    x: float
    y: float
    width: float
    height: float
    runs: list[TextRun] = field(default_factory=list)


@dataclass
class ChartData:
    chart_id: str
    chart_type: str
    title: str
    x: float = 0.0
    y: float = 0.0
    width: float = 0.0
    height: float = 0.0
    cache_values: list[list[str]] = field(default_factory=list)
    data_range_ref: Optional[str] = None
    source_citation: Optional[str] = None


@dataclass
class ImageElement:
    id: str
    x: float
    y: float
    width: float
    height: float
    image_data: Optional[bytes] = None
    content_type: Optional[str] = None
    extension: Optional[str] = None


@dataclass
class TableData:
    table_id: str
    rows: int
    columns: int
    title: str
    text: str
    x: float
    y: float
    width: float
    height: float


@dataclass
class SlideContent:
    slide_index: int
    title: str = ""
    text_boxes: list[TextBox] = field(default_factory=list)
    charts: list[ChartData] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    images: list[ImageElement] = field(default_factory=list)
    width: float = 10.0
    height: float = 7.5
    ocr_backend: str = "native"


@dataclass
class DeckContent:
    file_path: str
    slides: list[SlideContent] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    extracted_images: list[bytes] = field(default_factory=list)


@dataclass
class ExcelData:
    file_path: str
    sheets: dict[str, list[list[str]]] = field(default_factory=dict)
    named_ranges: dict[str, str] = field(default_factory=dict)


class DocumentIngestionService:
    def __init__(self):
        self.emu_to_inch = 914400
        self.pt_to_emu = 12700

    async def ingest_pptx(self, file_path: str) -> DeckContent:
        deck = DeckContent(file_path=file_path)

        prs = Presentation(file_path)
        deck.metadata["slide_count"] = len(prs.slides)
        deck.metadata["title"] = prs.core_properties.title or "Untitled"

        for idx, slide in enumerate(prs.slides):
            slide_content = self._extract_slide_content(slide, idx, prs)
            deck.slides.append(slide_content)
            deck.extracted_images.extend(
                img.image_data for img in slide_content.images if img.image_data
            )

        return deck

    async def ingest_pdf(self, file_path: str) -> DeckContent:
        import pdfplumber

        deck = DeckContent(file_path=file_path)

        try:
            with pdfplumber.open(file_path) as pdf:
                deck.metadata["slide_count"] = len(pdf.pages)
                deck.metadata["title"] = os.path.basename(file_path)

                # --- Pass 1: Extract text and rasterize pages needing OCR ---
                slides = []
                page_dims: dict[int, tuple[float, float]] = {}
                scanned_indices = []  # pages where text < 10 chars
                scanned_images = []  # rasterized PIL images for those pages

                # Check if OCR is needed for text-poor pages
                need_ocr = False
                try:
                    from ..agents.parallel_analysis import _should_skip_ocr_for_text_rich_pdf
                    if _should_skip_ocr_for_text_rich_pdf(file_path):
                        _get_doc_logger().info("PDF %s has embedded text; skipping OCR processing", file_path)
                        need_ocr = False
                    else:
                        need_ocr = True
                except Exception as e:
                    _get_doc_logger().warning("Failed to check OCR skip condition: %s (will attempt OCR)", e)
                    need_ocr = True

                for i, page in enumerate(pdf.pages):
                    page_w = float(page.width)
                    page_h = float(page.height)
                    text = page.extract_text() or ""

                    slide = SlideContent(
                        slide_index=i,
                        title=f"Page {i + 1}",
                        width=page_w / 72.0,
                        height=page_h / 72.0,
                    )
                    slides.append((slide, page, page_w, page_h, text))
                    page_dims[i] = (page_w, page_h)

                    if need_ocr and len(text.strip()) < 10:
                        img_rendered = None
                        try:
                            # 1. Try rendering with pypdfium2 (extremely fast, zero dependencies)
                            import pypdfium2 as pdfium
                            doc = pdfium.PdfDocument(file_path)
                            page_pdfium = doc[i]
                            # Render at 150 DPI (scale=2)
                            bitmap = page_pdfium.render(scale=2)
                            img_rendered = bitmap.to_pil()
                        except Exception as e_pdfium:
                            _get_doc_logger().warning("pdfium page render failed: %s", e_pdfium)
                            try:
                                img_rendered = page.to_image(resolution=150).original
                            except Exception as e_plumber:
                                _get_doc_logger().error("Both pdfium and pdfplumber page render failed: %s", e_plumber)
                        
                        scanned_images.append(img_rendered)
                        scanned_indices.append(i)

                # --- Pass 2: OCR scanned pages using multi-backend OCR pipeline ---
                ocr_texts: dict[int, str] = {}
                ocr_line_boxes: dict[int, list[dict]] = {}
                valid_imgs = [
                    (idx, img)
                    for idx, img in zip(scanned_indices, scanned_images)
                    if img is not None
                ]
                if valid_imgs:
                    try:
                        from app.services.ocr_asset_manager import OcrAssetManager
                        mgr = OcrAssetManager.get()
                        backend_id = mgr.active_backend_id()
                        from app.services.ocr_detectors import ocr_image
                        
                        for page_idx, pil_img in valid_imgs:
                            try:
                                ocr_result = ocr_image(pil_img, backend_id)
                                page_w, page_h = page_dims.get(page_idx, (0.0, 0.0))
                                img_w, img_h = pil_img.size
                                
                                if ocr_result and ocr_result.get("text"):
                                    ocr_texts[page_idx] = ocr_result["text"]
                                    
                                    line_boxes = []
                                    for line in (ocr_result.get("lines") or [])[:300]:
                                        bbox = line.get("bbox", {})
                                        x0 = bbox.get("x0", 0)
                                        y0 = bbox.get("y0", 0)
                                        x1 = bbox.get("x1", 1)
                                        y1 = bbox.get("y1", 1)
                                        line_text = (line.get("text") or "").strip()
                                        if not line_text:
                                            continue
                                        if img_w <= 0 or img_h <= 0 or page_w <= 0 or page_h <= 0:
                                            continue
                                        x_in = (float(x0) / float(img_w)) * (page_w / 72.0)
                                        y_in = (float(y0) / float(img_h)) * (page_h / 72.0)
                                        w_in = ((float(x1) - float(x0)) / float(img_w)) * (page_w / 72.0)
                                        h_in = ((float(y1) - float(y0)) / float(img_h)) * (page_h / 72.0)
                                        if w_in <= 0.02 or h_in <= 0.02:
                                            continue
                                        line_boxes.append({
                                            "text": line_text,
                                            "x": x_in,
                                            "y": y_in,
                                            "width": w_in,
                                            "height": h_in,
                                        })
                                    if line_boxes:
                                        ocr_line_boxes[page_idx] = line_boxes
                            except Exception as page_ocr_err:
                                _get_doc_logger().error("OCR failed for page %d: %s", page_idx, page_ocr_err)
                    except Exception as ocr_err:
                        _get_doc_logger().error("Multi-backend OCR pipeline failed: %s", ocr_err)

                # --- Pass 3: Build slides with resolved text ---
                for i, (slide, page, page_w, page_h, text) in enumerate(slides):
                    if i in ocr_texts and len(ocr_texts[i]) > len(text):
                        text = ocr_texts[i]
                        slide.title += " (OCR)"
                        slide.ocr_backend = "ocr"

                    line_boxes = ocr_line_boxes.get(i, [])
                    if line_boxes:
                        for box_idx, box in enumerate(line_boxes):
                            slide.text_boxes.append(
                                TextBox(
                                    id=f"tb_page_{i}_ocr_{box_idx}",
                                    text=box["text"],
                                    x=box["x"],
                                    y=box["y"],
                                    width=box["width"],
                                    height=box["height"],
                                )
                            )
                    else:
                        slide.text_boxes.append(
                            TextBox(
                                id=f"tb_page_{i}",
                                text=text,
                                x=0,
                                y=0,
                                width=page_w / 72.0,
                                height=page_h / 72.0,
                            )
                        )

                    # --- Extract tables WITH bounding boxes ---
                    found_tables = page.find_tables()
                    for j, tbl in enumerate(found_tables):
                        table_data_raw = tbl.extract()
                        if not table_data_raw:
                            continue
                        rows_count = len(table_data_raw)
                        cols_count = len(table_data_raw[0]) if rows_count > 0 else 0
                        table_text = "\n".join(
                            [
                                " | ".join([str(cell or "") for cell in row])
                                for row in table_data_raw
                            ]
                        )

                        # Convert bbox from PDF points to standard inches (72 points = 1 inch)
                        bbox = tbl.bbox  # (x0, top, x1, bottom) in PDF points
                        slide.tables.append(
                            TableData(
                                table_id=f"tbl_{i}_{j}",
                                rows=rows_count,
                                columns=cols_count,
                                title=f"Table {j + 1}",
                                text=table_text,
                                x=bbox[0] / 72.0,
                                y=bbox[1] / 72.0,
                                width=(bbox[2] - bbox[0]) / 72.0,
                                height=(bbox[3] - bbox[1]) / 72.0,
                            )
                        )

                    # --- Extract images from PDF page ---
                    try:
                        page_images = page.images or []
                        for j, img_info in enumerate(page_images):
                            img_x0 = float(img_info.get("x0", 0))
                            img_top = float(img_info.get("top", 0))
                            img_x1 = float(img_info.get("x1", img_x0))
                            img_bottom = float(img_info.get("bottom", img_top))

                            slide.images.append(
                                ImageElement(
                                    id=f"img_{i}_{j}",
                                    x=img_x0 / 72.0,
                                    y=img_top / 72.0,
                                    width=(img_x1 - img_x0) / 72.0,
                                    height=(img_bottom - img_top) / 72.0,
                                    image_data=None,  # Raw bytes extraction complex; rely on OCR backend vision
                                )
                            )
                    except Exception as img_err:
                        _get_doc_logger().error("PDF image extraction error on page %d: %s", i, img_err)

                    deck.slides.append(slide)
            return deck
        except Exception as e:
            error_str = str(e)
            if "PasswordIncorrect" in error_str or "encrypted" in error_str.lower():
                raise Exception(
                    "The uploaded PDF is encrypted. Please provide an unlocked PDF."
                )
            import logging

            logging.getLogger("slideforge.document_ingestion").error("Failed to ingest PDF: %s", error_str)
            raise Exception(f"Failed to ingest PDF: {error_str}")

    def _extract_slide_content(
        self, slide, idx: int, prs: Presentation
    ) -> SlideContent:
        content = SlideContent(slide_index=idx)

        slide_width = prs.slide_width / self.emu_to_inch
        slide_height = prs.slide_height / self.emu_to_inch
        content.width = slide_width
        content.height = slide_height

        title_text = ""
        for shape in self._iter_shapes(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if shape.name.lower().find("title") >= 0 and not title_text:
                    title_text = text
                if text:
                    tb = self._extract_text_box(shape, idx)
                    content.text_boxes.append(tb)

            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = self._extract_chart(shape, idx)
                if chart:
                    content.charts.append(chart)

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = self._extract_table(shape, idx)
                if table:
                    content.tables.append(table)

            # Extract images per-slide
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = self._extract_image(shape, idx)
                if img:
                    content.images.append(img)

        content.title = title_text or f"Slide {idx + 1}"
        return content

    def _iter_shapes(self, shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in self._iter_shapes(shape.shapes):
                        yield child
                except Exception:
                    continue
                continue
            yield shape

    def _extract_text_box(self, shape, slide_idx: int) -> TextBox:
        tb = TextBox(
            id=f"shape_{shape.shape_id}",
            text="",
            x=shape.left / self.emu_to_inch,
            y=shape.top / self.emu_to_inch,
            width=shape.width / self.emu_to_inch,
            height=shape.height / self.emu_to_inch,
        )

        if shape.has_text_frame:
            text_frame = shape.text_frame
            full_text = []
            char_offset = 0

            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text:
                        font = run.font
                        run_text = run.text
                        run_len = len(run_text)
                        text_run = TextRun(
                            text=run_text,
                            font_name=font.name or "Arial",
                            font_size=font.size.pt if font.size else 12,
                            font_bold=font.bold or False,
                            font_italic=font.italic or False,
                            font_color_rgb=self._get_font_color(font),
                            shape_id=str(shape.shape_id),
                            text_box_id=tb.id,
                            start_index=char_offset,
                            end_index=char_offset + run_len,
                        )
                        tb.runs.append(text_run)
                        full_text.append(run_text)
                        char_offset += run_len

            tb.text = "".join(full_text)

        return tb

    def _get_font_color(self, font) -> tuple[int, int, int]:
        try:
            if font.color and font.color.rgb:
                rgb = font.color.rgb
                if hasattr(rgb, "rgb"):
                    return (rgb.rgb[0], rgb.rgb[1], rgb.rgb[2])
                return (0, 0, 0)
        except Exception:
            pass
        return (0, 0, 0)

    def _extract_chart(self, shape, slide_idx: int) -> Optional[ChartData]:
        try:
            chart = shape.chart
            chart_data = ChartData(
                chart_id=f"chart_{shape.shape_id}",
                chart_type=str(chart.chart_type.name),
                title=chart.chart_title.text if chart.chart_title else "",
                x=shape.left / self.emu_to_inch,
                y=shape.top / self.emu_to_inch,
                width=shape.width / self.emu_to_inch,
                height=shape.height / self.emu_to_inch,
            )

            # Extract cache values from the chart's internal data
            try:
                for series in chart.series:
                    row = []
                    for point in series.values:
                        row.append(str(point) if point is not None else "")
                    if row:
                        chart_data.cache_values.append(row)
                if chart.series and hasattr(chart.series[0], "data_labels"):
                    chart_data.data_range_ref = str(
                        chart.series[0].values.address
                        if hasattr(chart.series[0].values, "address")
                        else ""
                    )
            except Exception:
                pass  # Cache extraction is best-effort; chart still returned without values

            return chart_data
        except Exception:
            return None

    def _extract_table(self, shape, slide_idx: int) -> Optional[TableData]:
        try:
            table = shape.table
            rows = len(table.rows)
            columns = len(table.columns) if rows > 0 else 0

            cell_texts = []
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        cell_texts.append(cell.text_frame.text.strip())

            table_data = TableData(
                table_id=f"table_{shape.shape_id}",
                rows=rows,
                columns=columns,
                title=f"Table {shape.shape_id}",
                text=" ".join(cell_texts),
                x=shape.left / self.emu_to_inch,
                y=shape.top / self.emu_to_inch,
                width=shape.width / self.emu_to_inch,
                height=shape.height / self.emu_to_inch,
            )

            return table_data
        except Exception:
            return None

    def _extract_image(self, shape, slide_idx: int) -> Optional[ImageElement]:
        try:
            img = ImageElement(
                id=f"img_{shape.shape_id}",
                x=shape.left / self.emu_to_inch,
                y=shape.top / self.emu_to_inch,
                width=shape.width / self.emu_to_inch,
                height=shape.height / self.emu_to_inch,
            )

            image_part = shape.image
            img.image_data = image_part.blob
            img.content_type = getattr(image_part, "content_type", None)
            img.extension = getattr(image_part, "ext", None)

            return img
        except Exception:
            return None

    async def ingest_excel(self, file_path: str) -> ExcelData:
        wb = load_workbook(file_path, data_only=True)

        excel_data = ExcelData(file_path=file_path)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):
                    rows.append(row_data)

            excel_data.sheets[sheet_name] = rows

        if wb.defined_names:
            for name, defn in wb.defined_names.items():
                excel_data.named_ranges[name] = str(defn.attr_text)

        return excel_data

    async def convert_pptx_to_images(
        self, file_path: str, output_dir: str, dpi: int = 150
    ) -> list[str]:
        output_paths = []
        os.makedirs(output_dir, exist_ok=True)

        rendered_via_com = False
        if os.name == "nt":
            try:
                import pythoncom
                import win32com.client

                pythoncom.CoInitialize()
                try:
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    abs_file_path = os.path.abspath(file_path)
                    abs_output_dir = os.path.abspath(output_dir)

                    # Open without showing window (minimize visual disruption)
                    presentation = powerpoint.Presentations.Open(
                        abs_file_path, WithWindow=False
                    )

                    # Export slides
                    for idx, slide in enumerate(presentation.Slides):
                        slide_image_path = os.path.join(
                            abs_output_dir, f"slide_{idx}.png"
                        )
                        slide.Export(slide_image_path, "PNG")
                        output_paths.append(slide_image_path)

                    presentation.Close()
                    powerpoint.Quit()
                    rendered_via_com = True
                finally:
                    pythoncom.CoUninitialize()
            except Exception as e:
                _get_doc_logger().warning("COM Rendering failed, falling back to placeholder: %s", e)

        if rendered_via_com:
            return output_paths

        # Fallback to placeholder rendering if COM fails or not on Windows
        prs = Presentation(file_path)

        for idx, slide in enumerate(prs.slides):
            slide_image_path = os.path.join(output_dir, f"slide_{idx}.png")

            try:
                # Note: True high-quality rendering usually requires COM objects on Windows
                # or LibreOffice/Uno. For this demo, we'll generate a better placeholder
                # that represents the slide structure.

                from PIL import Image, ImageDraw

                width_px = int(prs.slide_width / 12700)  # EMU to PX approx
                height_px = int(prs.slide_height / 12700)

                img = Image.new("RGB", (width_px, height_px), "white")
                draw = ImageDraw.Draw(img)

                # Draw a simple border
                draw.rectangle(
                    [0, 0, width_px - 1, height_px - 1], outline="#e2e8f0", width=2
                )

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if text:
                            x = int(shape.left / 12700)
                            y = int(shape.top / 12700)
                            w = int(shape.width / 12700)
                            h = int(shape.height / 12700)

                            # Draw box for text
                            draw.rectangle(
                                [x, y, x + w, y + h], outline="#f1f5f9", fill="#f8fafc"
                            )
                            draw.text(
                                (x + 5, y + 5),
                                text[:50] + ("..." if len(text) > 50 else ""),
                                fill="#475569",
                            )

                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        x = int(shape.left / 12700)
                        y = int(shape.top / 12700)
                        w = int(shape.width / 12700)
                        h = int(shape.height / 12700)
                        draw.rectangle(
                            [x, y, x + w, y + h], outline="#e2e8f0", fill="#eff6ff"
                        )
                        draw.text(
                            (x + w // 2 - 20, y + h // 2), "Image", fill="#3b82f6"
                        )

                img.save(slide_image_path)
                output_paths.append(slide_image_path)
            except Exception as e:
                _get_doc_logger().error("Error rendering slide %d: %s", idx, e)
                # Create a fallback blank image
                img = Image.new("RGB", (800, 600), "white")
                img.save(slide_image_path)
                output_paths.append(slide_image_path)

        return output_paths

    async def convert_pdf_to_images(
        self, file_path: str, output_dir: str, dpi: int = 150
    ) -> list[str]:
        output_paths = []
        os.makedirs(output_dir, exist_ok=True)

        # 1. Try with pypdfium2 (extremely fast, zero dependencies, pure python binding for Google PDFium)
        try:
            import pypdfium2 as pdfium
            _get_doc_logger().info("Converting PDF to images using pypdfium2...")
            doc = pdfium.PdfDocument(file_path)
            # Map DPI to rendering scale. 72 DPI is scale=1. 150 DPI is scale=2
            scale = max(1, int(dpi / 72))
            for i in range(len(doc)):
                image_path = os.path.join(output_dir, f"slide_{i}.png")
                page = doc[i]
                bitmap = page.render(scale=scale)
                pil_img = bitmap.to_pil()
                pil_img.save(image_path, format="PNG")
                output_paths.append(image_path)
            return output_paths
        except Exception as e:
            _get_doc_logger().warning("pypdfium2 conversion failed, falling back to pdfplumber: %s", e)

        # 2. Fallback to pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    image_path = os.path.join(output_dir, f"slide_{i}.png")
                    page.to_image(resolution=dpi).save(image_path, format="PNG")
                    output_paths.append(image_path)
            return output_paths
        except Exception as e:
            _get_doc_logger().error("Error converting PDF to images: %s", e)
            return []
