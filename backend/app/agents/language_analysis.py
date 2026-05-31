import re
import os
import asyncio
import logging
from dataclasses import dataclass, field

from ..models.schemas import Annotation
from ..services.llm_inference import inference_service, Message, parse_json_response


class LanguageToolError(Exception):
    pass


@dataclass
class LanguageResult:
    grammar_issues: list[Annotation] = field(default_factory=list)
    quality_issues: list[Annotation] = field(default_factory=list)
    tone_issues: list[Annotation] = field(default_factory=list)
    grammar_score: int = 100
    quality_score: int = 100
    tone_score: int = 100


HEDGING_PATTERNS = [
    r"\bcould potentially\b",
    r"\bmay wish to consider\b",
    r"\bthere is a possibility that\b",
    r"\bit might be\b",
    r"\bperhaps\b",
    r"\bmaybe\b",
    r"\bpossibly\b",
    r"\bto some extent\b",
    r"\bin some cases\b",
]

PASSIVE_PATTERNS = [
    r"\bwas\b.*\b( reduced| decreased| increased| declined| improved)\b",
    r"\bwere\b.*\b( conducted| performed| undertaken)\b",
    r"\bhas been\b.*\b( shown| demonstrated| observed)\b",
]

VAGUE_QUANTIFIERS = [
    r"\bsignificant\b",
    r"\bsubstantial\b",
    r"\bmeaningful\b",
    r"\bconsiderable\b",
    r"\bconsiderable\b",
    r"\bmajor\b",
    r"\bimportant\b",
]

RECOMMENDATION_NO_VERB_PATTERNS = [
    r"^(?!.*\b(?:manage|optimize|implement|divest|leverage|accelerate|reduce|increase|create|build|establish|launch|expand|transform)\b).*$",
]


from ..services.language_tool_client import LanguageToolClient  # noqa: E402


_logger_lang = logging.getLogger("slideforge.language_agent")


class LanguageAnalysisAgent:
    def __init__(self, language: str = "en-US"):
        self.language = language
        self.lt_client = LanguageToolClient()
        self._model = None  # Lazy load sentence-transformers
        self._max_concurrent_slides = int(
            os.environ.get("MAX_CONCURRENT_LANGUAGE", "6")
        )
        self._slide_semaphore = asyncio.Semaphore(max(1, self._max_concurrent_slides))

    async def analyze_slide(
        self, slide_index: int, text: str, guardrail_rules: dict = None
    ) -> LanguageResult:
        result = LanguageResult()

        grammar_issues = await self._check_grammar(text, slide_index)
        result.grammar_issues = grammar_issues
        result.grammar_score = self._calculate_grammar_score(grammar_issues)

        quality_issues = self._check_consulting_language(
            text, slide_index, guardrail_rules or {}
        )

        # F4 Task: Higher quality scoring via LLM
        llm_quality = await self._check_quality_llm(text, slide_index)
        quality_issues.extend(llm_quality)

        result.quality_issues = quality_issues
        result.quality_score = self._calculate_quality_score(quality_issues)

        # Tone issues are now handled at the deck level for consistency
        # result.tone_issues = await self._check_tone_llm(text, slide_index)
        # result.tone_score = 100 - (len(tone_issues) * 10)

        return result

    async def _check_quality_llm(self, text: str, slide_index: int) -> list[Annotation]:
        if not text or len(text.strip()) < 40 or inference_service.llm is None:
            return []

        prompt = f"""Analyze the following slide text for professional consulting quality (MBB standards).
Identify any issues with:
1. Lack of clarity or impact
2. "Fluff" or filler words
3. Weak action verbs

Text: "{text}"

Return a JSON list of issues, each with 'text' (the problematic snippet), 'message', and 'suggestion'.
If no issues, return [].
"""
        try:
            messages = [Message(role="user", content=prompt)]
            response = await inference_service.llm.generate(
                messages,
                max_tokens=min(
                    500, getattr(inference_service, "analysis_max_tokens", 500)
                ),
                context_window=getattr(inference_service, "local_context_window", None),
            )

            issues = parse_json_response(response.content)
            if not isinstance(issues, list):
                issues = []

            return [
                Annotation(
                    slide_index=slide_index,
                    text=issue.get("text", text[:20]) if isinstance(issue, dict) else text[:20],
                    category="quality",
                    severity="suggestion",
                    message=issue.get("message", "Professional quality issue") if isinstance(issue, dict) else "Professional quality issue",
                    suggestion=issue.get("suggestion") if isinstance(issue, dict) else None,
                )
                for issue in issues
                if isinstance(issue, dict)
            ]
        except Exception as e:
            _logger_lang.error("LLM quality check failed", exc_info=e)
            return []

    async def _check_tone_llm(self, text: str, slide_index: int) -> list[Annotation]:
        # Implementation for F4 Tone check
        prompt = f"Analyze the tone of this consulting slide. Is it objective and professional? Text: {text}. Return issues as JSON list [{{'text': '...', 'message': '...'}}] or []."
        try:
            messages = [Message(role="user", content=prompt)]
            response = await inference_service.llm.generate(
                messages,
                max_tokens=min(
                    300, getattr(inference_service, "analysis_max_tokens", 300)
                ),
                context_window=getattr(inference_service, "local_context_window", None),
            )
            issues = parse_json_response(response.content)
            if not isinstance(issues, list):
                issues = []
            return [
                Annotation(
                    slide_index=slide_index,
                    text=issue.get("text", "") if isinstance(issue, dict) else "",
                    category="tone",
                    severity="suggestion",
                    message=issue.get("message", "Tone improvement") if isinstance(issue, dict) else "Tone improvement",
                )
                for issue in issues
                if isinstance(issue, dict)
            ]
        except Exception:
            return []

    async def _check_grammar(self, text: str, slide_index: int) -> list[Annotation]:
        if not text or not text.strip():
            return []

        annotations = []

        # 1. External LanguageTool Check
        lt_matches = await self.lt_client.check(text, self.language)
        for match in lt_matches:
            rule_id = match.get("rule", {}).get("id", "")

            # Exclusion list
            if rule_id in ["WHITESPACE_RULE", "UPPERCASE_SENTENCE_START"]:
                continue

            severity = "warning"
            # Map typos to hard_blocks as per PRD
            if "MORFOLOGIK_RULE" in rule_id or "SPELLING" in rule_id:
                severity = "hard_block"

            annotations.append(
                Annotation(
                    slide_index=slide_index,
                    text=text[match["offset"] : match["offset"] + match["length"]],
                    category="grammar",
                    severity=severity,
                    message=match["message"],
                    suggestion=match["replacements"][0]["value"]
                    if match.get("replacements")
                    else None,
                )
            )

        # 2. Local Regex Fallback/Complement
        common_errors = {
            r"\btheir\b(?!\s+(is|are|was|were))": "Possible subject-verb agreement issue with 'their'",
            r"\bits\b(?=\s+[a-z]{3,})": "Potential article usage issue with 'its'",
            r",\s*,": "Double comma - remove one",
            r"\s+\.": "Extra space before period",
            r"\bdata\b(?=\s+(is|was|shows|suggests))": "'Data' is plural; consider 'data set' or 'information'",
        }

        for pattern, message in common_errors.items():
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                # Avoid duplicate annotations if LT already caught it
                if any(a.text == match.group() for a in annotations):
                    continue

                annotations.append(
                    Annotation(
                        slide_index=slide_index,
                        text=match.group(),
                        category="grammar",
                        severity="warning",
                        message=message,
                        suggestion=None,
                    )
                )

        return annotations

    def _check_consulting_language(
        self, text: str, slide_index: int, guardrail_rules: dict
    ) -> list[Annotation]:
        annotations = []

        prohibited = guardrail_rules.get("prohibited_phrases", [])
        for phrase in prohibited:
            if re.search(phrase, text, re.IGNORECASE):
                annotations.append(
                    Annotation(
                        slide_index=slide_index,
                        text=phrase,
                        category="prohibited_phrase",
                        severity="hard_block",
                        message=f"Prohibited phrase detected: {phrase}",
                        suggestion="Remove or rephrase this content",
                    )
                )

        for pattern in HEDGING_PATTERNS:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                annotations.append(
                    Annotation(
                        slide_index=slide_index,
                        text=match.group(),
                        category="hedging",
                        severity="warning",
                        message="Hedging language detected - consider being more direct",
                        suggestion=self._suggest_direct_alternative(match.group()),
                    )
                )

        for pattern in PASSIVE_PATTERNS:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                annotations.append(
                    Annotation(
                        slide_index=slide_index,
                        text=match.group(),
                        category="passive",
                        severity="warning",
                        message="Passive voice detected - consider specifying the actor",
                        suggestion="Add the responsible party for better clarity",
                    )
                )

        for pattern in VAGUE_QUANTIFIERS:
            if not re.search(r"\d+", text):
                matches = re.finditer(pattern, text, re.IGNORECASE)
                for match in matches:
                    if not any(
                        q in text.lower()
                        for q in ["percent", "%", "$", "million", "billion", "number"]
                    ):
                        annotations.append(
                            Annotation(
                                slide_index=slide_index,
                                text=match.group(),
                                category="vague_quantifier",
                                severity="warning",
                                message=f"Vague quantifier '{match.group()}' should be replaced with specific numbers",
                                suggestion="Add specific numbers or percentages",
                            )
                        )

        return annotations

    def _suggest_direct_alternative(self, hedging_phrase: str) -> str:
        alternatives = {
            "could potentially": "can",
            "may wish to consider": "should",
            "there is a possibility that": "it is likely that",
            "it might be": "it is",
            "perhaps": "likely",
            "maybe": "probably",
            "possibly": "likely",
        }

        for phrase, alt in alternatives.items():
            if phrase in hedging_phrase.lower():
                return alt

        return "use direct language"

    def _calculate_grammar_score(self, issues: list[Annotation]) -> int:
        if not issues:
            return 100

        deductions = len(issues) * 5
        return max(0, 100 - deductions)

    def _calculate_quality_score(self, issues: list[Annotation]) -> int:
        if not issues:
            return 100

        hard_blocks = sum(1 for i in issues if i.severity == "hard_block")
        warnings = sum(1 for i in issues if i.severity == "warning")

        deductions = hard_blocks * 15 + warnings * 5
        return max(0, 100 - deductions)

    async def analyze_deck(
        self, slides_data: list[dict], guardrail_rules: dict = None
    ) -> list[Annotation]:
        all_annotations = []

        # 1. Slide-by-slide Analysis — fire all slides concurrently
        async def _analyze_slide_limited(slide: dict):
            async with self._slide_semaphore:
                return await self.analyze_slide(
                    slide.get("index", 0),
                    slide.get("full_text", "") + " " + slide.get("title", ""),
                    guardrail_rules,
                )

        slide_tasks = [_analyze_slide_limited(slide) for slide in slides_data]
        slide_results = await asyncio.gather(*slide_tasks, return_exceptions=True)
        for result in slide_results:
            if isinstance(result, Exception):
                continue
            all_annotations.extend(result.grammar_issues)
            all_annotations.extend(result.quality_issues)
            all_annotations.extend(result.tone_issues)

        # 2. Global Deck-level Tone Consistency (GAP-04)
        tone_drift_annotations = await self.check_tone_consistency(slides_data)
        all_annotations.extend(tone_drift_annotations)

        return all_annotations

    async def check_tone_consistency(self, slides_data: list[dict]) -> list[Annotation]:
        """Uses sentence-transformers to detect slides with drifting tone."""
        if len(slides_data) < 3:
            return []

        try:
            from sentence_transformers import SentenceTransformer, util
            import torch

            if self._model is None:
                # Offline mode usually uses local paths or cached models
                self._model = SentenceTransformer("all-MiniLM-L6-v2")

            texts = [s.get("full_text", "") for s in slides_data]
            embeddings = self._model.encode(texts, convert_to_tensor=True)

            # Compute centroid embedding of the deck
            centroid = torch.mean(embeddings, dim=0)

            # Compute cosine similarities to centroid
            cosine_scores = util.cos_sim(embeddings, centroid.unsqueeze(0))

            annotations = []
            # Threshold for "drift" - usually 0.4-0.6 depending on model
            # Consultant slides should be very similar in tone
            threshold = 0.5

            drifted_slides = []
            for i, score in enumerate(cosine_scores[0]):
                if score < threshold:
                    drifted_slides.append(
                        {"index": slides_data[i].get("index", 0), "text": texts[i]}
                    )

            if not drifted_slides:
                return []

            # GAP-04: Perform single LLM call to summarize the tone drift across all flagged slides
            drift_summaries = await self._summarize_tone_drift(drifted_slides)

            for slide_index, summary in drift_summaries.items():
                annotations.append(
                    Annotation(
                        slide_index=slide_index,
                        text="[TONE DRIFT]",
                        category="tone",
                        severity="warning",
                        message=summary,
                        suggestion="Align phrasing and level of detail with previous slides.",
                    )
                )
            return annotations
        except Exception as e:
            logging.error(f"Tone consistency check failed: {e}")
            return []

    async def _summarize_tone_drift(self, drifted_slides: list[dict]) -> dict[int, str]:
        """Call LLM once to get human-readable drift descriptions for flagged slides."""
        if not drifted_slides:
            return {}

        prompt = "Analyze the tone of these consulting slides compared to the deck average (which is professional, objective, and consulting-style).\n\n"
        for slide in drifted_slides:
            prompt += f"--- Slide {slide['index']} ---\n{slide['text'][:500]}\n\n"

        prompt += "\nFor each slide index, provide a SHORT (1 sentence) description of WHY the tone is drifting (e.g. 'Too informal', 'Used sales-y language').\nReturn as JSON: { 'slide_index': 'reason' }"

        try:
            messages = [Message(role="user", content=prompt)]
            response = await inference_service.llm.generate(
                messages,
                max_tokens=min(
                    500, getattr(inference_service, "analysis_max_tokens", 500)
                ),
                context_window=getattr(inference_service, "local_context_window", None),
            )

            parsed = parse_json_response(response.content)
            if isinstance(parsed, dict):
                return {int(k): v for k, v in parsed.items()}
            return {}
        except Exception as e:
            _logger_lang.error("Tone drift summary failed", exc_info=e)
            return {
                s[
                    "index"
                ]: "Tone drift detected: language style differs from the rest of the deck."
                for s in drifted_slides
            }

    def write_comments_to_pptx(
        self, pptx_path: str, annotations: list[Annotation]
    ) -> str:
        from pptx import Presentation

        prs = Presentation(pptx_path)

        for annotation in annotations:
            if annotation.slide_index < len(prs.slides):
                slide = prs.slides[annotation.slide_index]

                comment_text = f"[{annotation.category}] {annotation.message}"
                if annotation.suggestion:
                    comment_text += f"\n\nSuggestion: {annotation.suggestion}"

                slide.comments.add_comment(
                    comment_text,
                    author_initials="SF",
                    author_name="SlideForge AI",
                )

        output_path = pptx_path.replace(".pptx", "_annotated.pptx")
        prs.save(output_path)

        return output_path
