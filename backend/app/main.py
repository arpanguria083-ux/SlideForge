from fastapi import (
    FastAPI,
    UploadFile,
    File,
    HTTPException,
    BackgroundTasks,
    Header,
    Depends,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import Optional
import os
import json
import re
import hashlib
import argparse
import zipfile
import shutil
import io
import logging
import threading
import time
import asyncio
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from contextlib import asynccontextmanager

from app.models.schemas import (
    GuardrailSchema,
    QAScorecard,
    Annotation,
    OverrideRequest,
    AcceptRequest,
    SignRequest,
    SaveTemplateRequest,
)
from app.services.document_ingestion import DocumentIngestionService
from app.services.guardrail import GuardrailManager
from app.services.adaptation_loop import adaptation_agent
from app.agents.parallel_analysis import (
    ParallelAnalysisOrchestrator,
    QAGradingOrchestrator,
    RevisionOrchestrator,
)
from app.agents.language_analysis import LanguageAnalysisAgent
from app.agents.template_discovery import template_discovery_agent
from app.services.audit_log import audit_log_service
from app.services.claim_evidence import ChromaDBManager, ClaimEvidenceGuardrail
from app.services.analysis_history import (
    AnalysisHistoryStore,
    ANALYSIS_HISTORY_VERSION,
)
from app.services.llm_inference import (
    InferenceProvider,
    inference_service,
    LLMFactory,
    Message,
    parse_json_response,
)
from app.services.model_registry import model_registry
from app.core.session_store import SQLiteSessionStore
from app.api import api_router
from app.core.config import AppSettings
from app.core.request_id import request_id_middleware
from app.core.time_utils import utc_now_compact, utc_now_iso


# Configure structured logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(name)s] %(levelname)s: %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("slideforge")


class RedactingFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:
        rendered = str(record.getMessage())
        lowered = rendered.lower()
        if "authorization" in lowered or "api_key" in lowered:
            record.msg = "[REDACTED SENSITIVE LOG CONTENT]"
            record.args = ()
        return True


for handler in logging.getLogger().handlers:
    handler.addFilter(RedactingFilter())

settings = AppSettings()


@asynccontextmanager
async def _lifespan(_app: FastAPI):
    await startup_event()
    try:
        yield
    finally:
        await shutdown_event()


app = FastAPI(title=settings.app_name, lifespan=_lifespan)
app.include_router(api_router)
app.middleware("http")(request_id_middleware)
_service_load_lock = threading.Lock()


class RuntimeProviderConfigPayload(BaseModel):
    provider: str
    api_base_url: Optional[str] = None
    api_key: Optional[str] = None
    model: Optional[str] = None
    local_context_window: Optional[int] = None


def _normalize_base_url(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    normalized = value.strip()
    if not normalized:
        return ""
    return normalized.rstrip("/")


def _mask_secret(value: str) -> Optional[str]:
    if not value:
        return None
    if len(value) <= 6:
        return "*" * len(value)
    return f"{value[:4]}{'*' * max(4, len(value) - 6)}{value[-2:]}"


def _provider_requires_api_key(provider: InferenceProvider) -> bool:
    return provider == InferenceProvider.API


def _provider_is_configured(provider: InferenceProvider) -> bool:
    try:
        return inference_service.provider_is_ready(provider)
    except ValueError:
        return (
            inference_service.current_provider == provider.value
            and inference_service.llm is not None
        )


def _provider_response(provider: InferenceProvider) -> dict:
    try:
        config = inference_service.get_provider_connection_config(provider)
        base_url = config.base_url
        model = config.model
        api_key = config.api_key
    except ValueError:
        base_url = ""
        model = provider.value
        api_key = ""
    return {
        "api_base_url": base_url,
        "configured": _provider_is_configured(provider),
        "requires_api_key": _provider_requires_api_key(provider),
        "base_url": base_url,
        "model": model,
        "api_key_configured": bool(api_key),
        "api_key_preview": _mask_secret(api_key),
    }


def _runtime_provider_config_response() -> dict:
    try:
        current_provider = InferenceProvider(inference_service.current_provider)
    except ValueError:
        current_provider = InferenceProvider.API
    current_config = _provider_response(current_provider)
    return {
        "enabled": True,
        "provider": current_provider.value,
        "api_base_url": current_config["api_base_url"],
        "model": current_config["model"],
        "local_context_window": inference_service.local_context_window,
        "api_key_configured": current_config["api_key_configured"],
        "api_key_preview": current_config["api_key_preview"],
        "configured": current_config["configured"],
        "requires_api_key": current_config["requires_api_key"],
        "providers": {
            InferenceProvider.API.value: _provider_response(InferenceProvider.API),
            InferenceProvider.MLX.value: _provider_response(InferenceProvider.MLX),
            InferenceProvider.TRANSFORMERS.value: _provider_response(
                InferenceProvider.TRANSFORMERS
            ),
            InferenceProvider.OLLAMA.value: _provider_response(
                InferenceProvider.OLLAMA
            ),
            InferenceProvider.LM_STUDIO.value: _provider_response(
                InferenceProvider.LM_STUDIO
            ),
        },
    }


def _extract_provider_updates(
    payload: RuntimeProviderConfigPayload,
) -> Optional[dict]:
    fields_set = getattr(payload, "model_fields_set", set())
    updates = {}
    if "api_base_url" in fields_set:
        updates["base_url"] = _normalize_base_url(payload.api_base_url)
    if "api_key" in fields_set:
        updates["api_key"] = (payload.api_key or "").strip()
    if "model" in fields_set:
        updates["model"] = (payload.model or "").strip()
    return updates or None


def _enforce_upload_size(content: bytes) -> None:
    max_size = int(settings.max_file_size)
    if len(content) > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"Uploaded file exceeds max size of {max_size} bytes",
        )


def ensure_services_loaded():
    global ingestion_service, guardrail_manager, analysis_orchestrator
    global qa_grader, revision_orchestrator, language_agent
    global chroma_manager, claim_evidence_guardrail, analysis_history_store

    if all(
        service is not None
        for service in (
            ingestion_service,
            guardrail_manager,
            analysis_orchestrator,
            qa_grader,
            revision_orchestrator,
            language_agent,
            chroma_manager,
            claim_evidence_guardrail,
            analysis_history_store,
        )
    ):
        return

    with _service_load_lock:
        if all(
            service is not None
            for service in (
                ingestion_service,
                guardrail_manager,
                analysis_orchestrator,
                qa_grader,
                revision_orchestrator,
                language_agent,
                chroma_manager,
                claim_evidence_guardrail,
                analysis_history_store,
            )
        ):
            return

        ingestion_service = DocumentIngestionService()
        guardrail_manager = GuardrailManager(str(data_dir / "guardrails"))
        analysis_orchestrator = ParallelAnalysisOrchestrator()
        qa_grader = QAGradingOrchestrator()
        revision_orchestrator = RevisionOrchestrator()
        language_agent = LanguageAnalysisAgent()

        chroma_manager = ChromaDBManager(str(data_dir / "chromadb"))
        claim_evidence_guardrail = ClaimEvidenceGuardrail(chroma_manager)
        analysis_history_store = AnalysisHistoryStore(
            str(data_dir / "analysis_history")
        )
        logging.info("Lazy-loaded heavy ML models and services.")


async def startup_event():
    recovered = session_store.mark_stale_inflight_failed(stale_seconds=60)
    if recovered:
        logger.warning("Recovered %s stale in-flight sessions as failed", recovered)
    active_sessions.update(session_store.list_active(_session_ttl_seconds()))
    asyncio.create_task(asyncio.to_thread(ensure_services_loaded))
    asyncio.create_task(_cleanup_expired_sessions_loop())
    from app.core.preflight import run_preflight_checks

    results = run_preflight_checks(data_dir=str(settings.data_dir))
    for check in results["checks"]:
        if check["status"] != "OK":
            logging.warning(f"[PREFLIGHT] {check['name']}: {check['message']}")


cors_origins = os.environ.get(
    "CORS_ORIGINS", "http://127.0.0.1:3000,http://localhost:3000"
).split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=cors_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

data_dir = Path(settings.data_dir)
data_dir.mkdir(parents=True, exist_ok=True)

ingestion_service = None
guardrail_manager = None
analysis_orchestrator = None
qa_grader = None
revision_orchestrator = None
language_agent = None
chroma_manager = None
claim_evidence_guardrail = None
analysis_history_store = None


active_sessions = {}
session_store = SQLiteSessionStore(str(data_dir / "sessions.db"))
_cleanup_stats_lock = threading.Lock()
_cleanup_stats = {
    "last_run_ts": None,
    "last_expired_count": 0,
    "last_expired_ids": [],
    "total_expired_count": 0,
}


def _touch_session(session: dict):
    session["last_access_ts"] = time.time()


def _persist_session(session_id: str, session: dict):
    try:
        session_store.save(session_id, session)
    except Exception:
        logger.exception("Failed to persist session %s", session_id)


def _namespace_key(client_namespace: Optional[str]) -> str:
    cleaned = (client_namespace or "").strip()
    return cleaned if cleaned else "__anonymous__"


def _session_ttl_seconds() -> int:
    return max(3600, int(settings.session_ttl_hours * 3600))


def _session_last_access_ts(session: dict, now_ts: float) -> float:
    return float(session.get("last_access_ts", session.get("created_at_ts", now_ts)))


def _is_session_expired(session: dict, now_ts: float, ttl_seconds: int) -> bool:
    last_access_ts = _session_last_access_ts(session, now_ts)
    return now_ts - last_access_ts > ttl_seconds


def _get_session_or_404(session_id: str) -> dict:
    session = active_sessions.get(session_id)
    if session is None:
        restored = session_store.load(session_id)
        if restored is not None:
            active_sessions[session_id] = restored
            session = restored
    if session is None:
        raise HTTPException(status_code=404, detail="Session not found")
    _touch_session(session)
    _persist_session(session_id, session)
    return session


def _cleanup_session_artifacts(session_id: str, session: dict):
    upload_root = (data_dir / "uploads").resolve()
    sessions_root = (data_dir / "sessions").resolve()
    candidates = [upload_root / session_id, sessions_root / session_id]

    for candidate in candidates:
        try:
            resolved = candidate.resolve()
            if resolved.exists() and (
                str(resolved).startswith(str(upload_root))
                or str(resolved).startswith(str(sessions_root))
            ):
                shutil.rmtree(resolved, ignore_errors=True)
        except Exception:
            logger.exception("Failed to clean up session path: %s", candidate)

    source_namespace = session.get("source_namespace")
    if source_namespace and chroma_manager is not None:
        try:
            chroma_manager.delete_collection(source_namespace)
        except Exception:
            logger.exception(
                "Failed to clean up Chroma namespace for session %s", session_id
            )


def _extract_docx_text(content: bytes) -> list[str]:
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        try:
            xml_bytes = archive.read("word/document.xml")
        except KeyError as exc:
            raise ValueError("DOCX file is missing word/document.xml") from exc

    root = ET.fromstring(xml_bytes)
    for paragraph in root.findall(".//w:p", namespace):
        runs = [
            node.text for node in paragraph.findall(".//w:t", namespace) if node.text
        ]
        text = "".join(runs).strip()
        if text:
            paragraphs.append(text)

    return ["\n".join(paragraphs)] if paragraphs else []


def _extract_uploaded_document_text(
    filename: str, content: bytes, filepath: Optional[Path] = None
) -> list[str]:
    ext = os.path.splitext(filename)[1].lower()
    if ext in {".txt", ".csv", ".md"}:
        return [content.decode("utf-8", errors="ignore")]

    if ext == ".pdf":
        import pdfplumber

        pdf_source = str(filepath) if filepath else io.BytesIO(content)
        documents: list[str] = []
        with pdfplumber.open(pdf_source) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    documents.append(text)
        return documents

    if ext in {".xlsx", ".xls"}:
        from openpyxl import load_workbook

        workbook_source = str(filepath) if filepath else io.BytesIO(content)
        wb = load_workbook(workbook_source, data_only=True)
        documents: list[str] = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                rows.append(" | ".join([str(cell or "") for cell in row]))
            if rows:
                documents.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        return documents

    if ext == ".docx":
        return _extract_docx_text(content)

    raise ValueError(f"Unsupported document extension: {ext}")


def _record_cleanup_stats(expired_ids: list[str]):
    now_ts = time.time()
    with _cleanup_stats_lock:
        _cleanup_stats["last_run_ts"] = now_ts
        _cleanup_stats["last_expired_count"] = len(expired_ids)
        _cleanup_stats["last_expired_ids"] = expired_ids[:50]
        _cleanup_stats["total_expired_count"] = int(
            _cleanup_stats.get("total_expired_count", 0)
        ) + len(expired_ids)


def _collect_expired_session_ids(now_ts: Optional[float] = None) -> list[str]:
    now = now_ts if now_ts is not None else time.time()
    ttl_seconds = _session_ttl_seconds()
    expired_ids = []
    for sid, session in list(active_sessions.items()):
        if _is_session_expired(session, now, ttl_seconds):
            expired_ids.append(sid)
    return expired_ids


def _evict_expired_sessions_now(now_ts: Optional[float] = None) -> list[str]:
    expired_ids = _collect_expired_session_ids(now_ts)
    if not expired_ids:
        _record_cleanup_stats([])
        return []

    for sid in expired_ids:
        session = active_sessions.pop(sid, None)
        if session is not None:
            _cleanup_session_artifacts(sid, session)
        session_store.delete(sid)

    _record_cleanup_stats(expired_ids)
    logger.info("Expired %s inactive session(s): %s", len(expired_ids), expired_ids)
    return expired_ids


def _count_sessions_by_namespace() -> dict[str, int]:
    counts: dict[str, int] = {}
    for session in active_sessions.values():
        ns_key = session.get("namespace_key") or _namespace_key(
            session.get("client_namespace")
        )
        counts[ns_key] = counts.get(ns_key, 0) + 1
    return counts


def _enforce_session_capacity(namespace_key: str):
    _evict_expired_sessions_now()

    max_total = max(1, int(settings.max_active_sessions))
    if len(active_sessions) >= max_total:
        raise HTTPException(
            status_code=429,
            detail=f"Active session limit reached ({max_total}). Retry later.",
        )

    max_per_ns = int(settings.max_active_sessions_per_namespace)
    if max_per_ns > 0:
        by_namespace = _count_sessions_by_namespace()
        if by_namespace.get(namespace_key, 0) >= max_per_ns:
            raise HTTPException(
                status_code=429,
                detail=(
                    f"Active session limit reached for namespace '{namespace_key}' "
                    f"({max_per_ns}). Retry later."
                ),
            )


async def _cleanup_expired_sessions_loop():
    interval_seconds = max(60, int(settings.session_cleanup_interval_minutes * 60))

    while True:
        await asyncio.sleep(interval_seconds)
        _evict_expired_sessions_now()


def _clamp_percent(value: float) -> float:
    return max(0.0, min(100.0, value))


def _deterministic_band_index(text: str, modulo: int = 6) -> int:
    if modulo <= 1:
        return 0
    digest = hashlib.sha1((text or "").encode("utf-8", errors="ignore")).hexdigest()
    return int(digest[:8], 16) % modulo


def _build_annotation_id(annotation: dict) -> str:
    slide_index = annotation.get("slide_index", 0)
    payload = {
        "message": annotation.get("message") or "",
        "text": annotation.get("text") or "",
        "category": annotation.get("category") or "",
        "severity": annotation.get("severity") or "",
        "shape_id": annotation.get("shape_id") or "",
        "run_start": annotation.get("run_start"),
        "run_end": annotation.get("run_end"),
    }
    digest = hashlib.sha1(
        json.dumps(payload, sort_keys=True).encode("utf-8", errors="ignore")
    ).hexdigest()[:12]
    return f"{slide_index}:{digest}"


def _group_annotations_by_slide(annotations: list[dict]) -> dict[str, list[dict]]:
    grouped: dict[str, list[dict]] = {}
    for annotation in annotations:
        key = str(annotation.get("slide_index", 0))
        grouped.setdefault(key, []).append(annotation)
    return grouped


def _serialize_annotation(annotation) -> dict:
    if hasattr(annotation, "model_dump"):
        return annotation.model_dump()
    return dict(annotation)


def _clone_slide_payloads_for_session(
    session_id: str, slides_data: list[dict], upload_dir: Path
) -> list[dict]:
    cloned: list[dict] = []
    for slide in slides_data:
        slide_copy = json.loads(json.dumps(slide))
        slide_index = slide_copy.get("index", 0)
        slide_copy["preview_path"] = str(
            upload_dir / "previews" / f"slide_{slide_index}.png"
        )
        slide_copy["previewUrl"] = (
            f"/api/session/{session_id}/slide/{slide_index}/image"
        )
        for image in slide_copy.get("images", []) or []:
            asset_path = image.get("asset_path")
            if asset_path and Path(asset_path).exists():
                image["asset_url"] = (
                    f"/api/session/{session_id}/slide/{slide_index}/asset/{image.get('id', '')}"
                )
            else:
                image["asset_path"] = None
                image["asset_url"] = None

        slide_copy.setdefault("ocr_backend", "native")
        cloned.append(slide_copy)
    return cloned


def _history_entry_is_current(history_entry: Optional[dict]) -> bool:
    if not history_entry:
        return False
    if int(history_entry.get("analysis_version", 0)) < ANALYSIS_HISTORY_VERSION:
        return False
    deep_analysis = history_entry.get("deep_analysis_by_slide", {})
    if not deep_analysis:
        return False
    first_item = next(iter(deep_analysis.values()), {})
    return bool(first_item.get("review"))


async def _generate_previews_for_deck(deck_path: str, upload_dir: Path):
    ensure_services_loaded()
    previews_dir = upload_dir / "previews"
    previews_dir.mkdir(parents=True, exist_ok=True)
    if deck_path.lower().endswith(".pptx"):
        await ingestion_service.convert_pptx_to_images(deck_path, str(previews_dir))
    elif deck_path.lower().endswith(".pdf"):
        await ingestion_service.convert_pdf_to_images(deck_path, str(previews_dir))


def _slides_missing_browser_assets(slides_data: list[dict]) -> bool:
    for slide in slides_data or []:
        for image in slide.get("images", []) or []:
            if image.get("has_content") and not image.get("asset_url"):
                return True
    return False


async def _hydrate_slide_assets_for_session(session_id: str, session: dict):
    deck_path = session.get("deck_path")
    slides_data = session.get("slides_data", [])
    if not deck_path or not str(deck_path).lower().endswith(".pptx") or not slides_data:
        return
    if not _slides_missing_browser_assets(slides_data):
        return

    ensure_services_loaded()
    deck_content = await ingestion_service.ingest_pptx(deck_path)
    upload_dir = Path(deck_path).parent
    assets_dir = upload_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    slides_by_index = {slide.get("index", 0): slide for slide in slides_data}

    for parsed_slide in deck_content.slides:
        session_slide = slides_by_index.get(parsed_slide.slide_index)
        if not session_slide:
            continue

        existing_images = session_slide.get("images", []) or []
        hydrated_images = []
        for idx, img in enumerate(parsed_slide.images):
            existing = existing_images[idx] if idx < len(existing_images) else {}
            asset_path = None
            asset_url = None
            extension = (img.extension or "bin").lower().lstrip(".")
            content_type = img.content_type or "application/octet-stream"
            if img.image_data:
                asset_filename = (
                    f"slide_{parsed_slide.slide_index}_{img.id}.{extension}"
                )
                asset_file = assets_dir / asset_filename
                asset_file.write_bytes(img.image_data)
                asset_path = str(asset_file)
                asset_url = f"/api/session/{session_id}/slide/{parsed_slide.slide_index}/asset/{img.id}"

            hydrated_images.append(
                {
                    **existing,
                    "id": img.id,
                    "x": img.x,
                    "y": img.y,
                    "width": img.width,
                    "height": img.height,
                    "has_content": img.image_data is not None,
                    "coord_unit": existing.get("coord_unit") or "absolute",
                    "asset_path": asset_path,
                    "asset_url": asset_url,
                    "content_type": content_type,
                    "extension": extension if img.image_data else None,
                }
            )

        if hydrated_images:
            session_slide["images"] = hydrated_images


async def _restore_history_to_session(
    session_id: str, session: dict, history_entry: dict
):
    archived_deck_path = history_entry.get("archived_deck_path")
    if not archived_deck_path or not Path(archived_deck_path).exists():
        raise HTTPException(
            status_code=404, detail="Archived deck not found for this history entry"
        )

    upload_dir = data_dir / "uploads" / session_id
    upload_dir.mkdir(parents=True, exist_ok=True)
    dest_deck_path = upload_dir / Path(archived_deck_path).name
    shutil.copy2(archived_deck_path, dest_deck_path)
    await _generate_previews_for_deck(str(dest_deck_path), upload_dir)

    slides_data = _clone_slide_payloads_for_session(
        session_id, history_entry.get("slides_data", []), upload_dir
    )

    session.update(
        {
            "deck_path": str(dest_deck_path),
            "document_fingerprint": history_entry.get("fingerprint"),
            "original_filename": history_entry.get("original_filename"),
            "deck_metadata": history_entry.get("deck_metadata", {}),
            "slides_data": slides_data,
            "source_namespace": None,
            "source_files": [],
            "source_indexed_chunks": {},
            "excel_data": None,
            "scorecard": history_entry.get("scorecard", {}),
            "annotations_by_slide": _group_annotations_by_slide(
                history_entry.get("scorecard", {}).get("annotations", [])
            ),
            "agent_metadata": history_entry.get("agent_metadata", {}),
            "deep_analysis_by_slide": history_entry.get("deep_analysis_by_slide", {}),
            "status": "analyzed",
            "history_restored": True,
            "history_restored_at": utc_now_iso(),
        }
    )


def _build_deep_analysis_by_slide(
    agent_results: list, language_annotations: list, slides_data: list[dict]
) -> dict[str, dict]:
    deep_by_slide: dict[str, dict] = {}

    for slide in slides_data:
        slide_idx = slide.get("index", 0)
        deep_by_slide[str(slide_idx)] = {
            "agents": [],
            "judge": {"name": "Language Analysis", "findings": [], "summary": ""},
        }

    for result in agent_results:
        grouped = _group_annotations_by_slide(
            [_serialize_annotation(a) for a in result.findings]
        )
        for slide_idx, payload in deep_by_slide.items():
            findings = grouped.get(slide_idx, [])
            payload["agents"].append(
                {
                    "name": result.agent_name,
                    "score": result.score,
                    "findings": findings,
                    "summary": (
                        (
                            "No slide-specific issues detected, but this agent found issues elsewhere in the deck."
                            if not findings and result.score < 85
                            else "No major issues detected by this agent."
                        )
                        if not findings
                        else f"{len(findings)} finding(s) on this slide."
                    ),
                }
            )

    language_grouped = _group_annotations_by_slide(
        [_serialize_annotation(a) for a in language_annotations]
    )
    for slide_idx, payload in deep_by_slide.items():
        findings = language_grouped.get(slide_idx, [])
        payload["judge"] = {
            "name": "Language Analysis",
            "findings": findings,
            "summary": (
                "Language and tone checks passed."
                if not findings
                else f"{len(findings)} language or tone issue(s) detected."
            ),
        }

    return deep_by_slide


# ---------------------------------------------------------------------------
# Scoring functions — delegated to services.scoring module
# ---------------------------------------------------------------------------
from .services.scoring import (
    build_slide_consultant_score as _build_slide_consultant_score,
    build_slide_reliability as _build_slide_reliability,
)


def _weighted_dimension_score(
    scores: dict[str, float], weights: dict[str, float]
) -> int:
    weighted = sum(float(scores.get(k, 0)) * float(weights.get(k, 0)) for k in weights)
    return max(0, min(100, int(round(weighted))))


def _normalize_guardrail_weights(guardrail: GuardrailSchema) -> dict[str, float]:
    return guardrail.normalized_rubric_weights()


def _invalidate_session_analysis(session: dict) -> None:
    session.pop("scorecard", None)
    session.pop("annotations_by_slide", None)
    session.pop("deep_analysis_by_slide", None)
    session.pop("agent_metadata", None)
    session.pop("auto_fixes", None)
    session["history_restored"] = False
    session["status"] = "parsed"


def _session_id_by_ref(session: dict) -> Optional[str]:
    for sid, current in active_sessions.items():
        if current is session:
            return sid
    return None


def _persist_known_session(session: dict) -> None:
    sid = _session_id_by_ref(session)
    if sid:
        _persist_session(sid, session)


def _refresh_scorecard_counters(session: dict) -> None:
    scorecard_data = session.get("scorecard")
    if not isinstance(scorecard_data, dict):
        return
    annotations = scorecard_data.get("annotations", []) or []
    hard_blocks = sum(1 for a in annotations if a.get("severity") == "hard_block")
    warnings = sum(1 for a in annotations if a.get("severity") == "warning")
    failing_slides = sorted(
        {
            int(a.get("slide_index", 0))
            for a in annotations
            if isinstance(a.get("slide_index", 0), int)
        }
    )
    scorecard_data["hard_block_count"] = hard_blocks
    scorecard_data["warning_count"] = warnings
    scorecard_data["failing_slides"] = failing_slides
    session["scorecard"] = scorecard_data


def _remove_annotation_from_deep_analysis(session: dict, annotation_id: str) -> None:
    deep = session.get("deep_analysis_by_slide")
    if not isinstance(deep, dict):
        return

    for payload in deep.values():
        if not isinstance(payload, dict):
            continue

        judge = payload.get("judge")
        if isinstance(judge, dict):
            findings = judge.get("findings", []) or []
            next_findings = []
            for item in findings:
                if not isinstance(item, dict):
                    next_findings.append(item)
                    continue
                candidate_ids = {
                    _build_annotation_id(item),
                    item.get("message", ""),
                    item.get("text", ""),
                }
                if annotation_id in candidate_ids:
                    continue
                next_findings.append(item)
            judge["findings"] = next_findings

        for agent in payload.get("agents", []) or []:
            if not isinstance(agent, dict):
                continue
            findings = agent.get("findings", []) or []
            next_findings = []
            for item in findings:
                if not isinstance(item, dict):
                    next_findings.append(item)
                    continue
                candidate_ids = {
                    _build_annotation_id(item),
                    item.get("message", ""),
                    item.get("text", ""),
                }
                if annotation_id in candidate_ids:
                    continue
                next_findings.append(item)
            agent["findings"] = next_findings

    session["deep_analysis_by_slide"] = deep


def _build_dynamic_slide_scorecard(
    guardrail: GuardrailSchema,
    consultant_score: dict,
    reliability: dict,
    slide_annotations: list[dict],
    guardrail_coverage: list[dict],
    framework_meta: dict,
    so_what_meta: dict,
    benchmark_meta: dict,
) -> dict:
    weights = _normalize_guardrail_weights(guardrail)
    breakdown = consultant_score.get("breakdown", {}) or {}

    msg = float(breakdown.get("message_clarity", 0) or 0)
    evd = float(breakdown.get("evidence_strength", 0) or 0)
    layout = float(breakdown.get("layout_quality", 0) or 0)
    visual = float(breakdown.get("visual_usefulness", 0) or 0)

    framework_score = float(framework_meta.get("score", msg) or msg)
    so_what_score = float(so_what_meta.get("score", msg) or msg)
    benchmark_score = float(benchmark_meta.get("score", evd) or evd)

    dimension_scores: dict[str, float] = {
        "structure": round(max(0.0, min(100.0, (msg * 0.7 + layout * 0.3))), 2),
        "claim_grounding": round(max(0.0, min(100.0, evd)), 2),
        "data_accuracy": round(max(0.0, min(100.0, evd)), 2),
        "visual": round(max(0.0, min(100.0, visual)), 2),
        "language": round(max(0.0, min(100.0, msg)), 2),
        "framework": round(max(0.0, min(100.0, framework_score)), 2),
        "so_what": round(max(0.0, min(100.0, so_what_score)), 2),
        "benchmarking": round(max(0.0, min(100.0, benchmark_score)), 2),
    }

    overall = _weighted_dimension_score(dimension_scores, weights)
    hard_blocks = sum(1 for a in slide_annotations if a.get("severity") == "hard_block")
    warnings = sum(1 for a in slide_annotations if a.get("severity") == "warning")

    pass_threshold = int(guardrail.pass_threshold or 75)
    warn_threshold = max(0, pass_threshold - 15)
    if hard_blocks > 0:
        status = "hard_fail"
    elif overall >= pass_threshold:
        status = "pass"
    elif overall >= warn_threshold:
        status = "warn"
    else:
        status = "fail"

    checked = sum(1 for item in guardrail_coverage if item.get("status") == "checked")
    failed = sum(1 for item in guardrail_coverage if item.get("status") == "failed")
    skipped = sum(1 for item in guardrail_coverage if item.get("status") == "skipped")

    dimensions = []
    for dim, score in dimension_scores.items():
        dim_status = "checked" if score >= pass_threshold else "failed"
        dimensions.append(
            {
                "id": dim,
                "score": int(round(score)),
                "weight": round(float(weights.get(dim, 0.0)), 6),
                "status": dim_status,
            }
        )

    return {
        "schema_version": "1.0",
        "weights": weights,
        "thresholds": {
            "pass": pass_threshold,
            "warn": warn_threshold,
        },
        "dimensions": dimensions,
        "coverage": {
            "rules_total": len(guardrail_coverage),
            "rules_checked": checked,
            "rules_failed": failed,
            "rules_skipped": skipped,
        },
        "overall": {
            "score": overall,
            "status": status,
            "confidence_score": int(reliability.get("score", 0) or 0),
            "hard_block_count": hard_blocks,
            "warning_count": warnings,
        },
    }


async def _build_source_grounding_annotations(session: dict) -> list[Annotation]:
    source_namespace = session.get("source_namespace")
    slides_data = session.get("slides_data", [])
    if not source_namespace or claim_evidence_guardrail is None:
        return []

    session_guardrail = _guardrail_from_session(session)
    citation_rules = session_guardrail.language_rules if session_guardrail else {}

    grounding_annotations: list[Annotation] = []
    for slide in slides_data:
        slide_index = slide.get("index", 0)
        full_text = (slide.get("full_text") or "").strip()
        if len(full_text) < 20:
            continue

        try:
            if inference_service.llm is not None:
                claims = await inference_service.extract_claims(
                    full_text, citation_rules
                )
            else:
                claims = []
        except Exception:
            claims = []

        if not claims:
            heuristic_claims = []
            segments = re.split(r"(?<=[.!?])\s+|\n+", full_text)
            for segment in segments:
                cleaned = segment.strip()
                if len(cleaned) < 30:
                    continue
                if re.search(r"\d", cleaned) or any(
                    token in cleaned.lower()
                    for token in (
                        "according to",
                        "increase",
                        "decrease",
                        "market",
                        "revenue",
                        "cost",
                        "growth",
                        "benchmark",
                    )
                ):
                    heuristic_claims.append({"claim": cleaned, "confidence": 0.45})
                if len(heuristic_claims) >= 8:
                    break
            claims = heuristic_claims

        claim_texts = []
        for claim in claims:
            if isinstance(claim, dict):
                claim_text = (claim.get("claim") or "").strip()
                if claim_text:
                    claim_texts.append(claim_text)

        if not claim_texts:
            continue

        try:
            grounding_results = await claim_evidence_guardrail.check_claims_batch(
                claim_texts, source_namespace
            )
        except Exception:
            continue

        for result in grounding_results:
            if not isinstance(result, dict) or result.get("grounded", False):
                continue

            grounding_annotations.append(
                Annotation(
                    slide_index=slide_index,
                    text=result.get("claim", full_text[:80]),
                    category="claim_grounding",
                    severity=result.get("severity", "hard_block"),
                    message=result.get(
                        "message", "Claim is not grounded in uploaded source evidence"
                    ),
                    suggestion=(
                        f"Add support from uploaded sources or revise the claim. Evidence checked: "
                        f"{(result.get('evidence') or '')[:160]}"
                    ).strip(),
                )
            )

    return grounding_annotations


async def _enrich_deep_analysis_with_slide_reviews(
    slides_data: list[dict], deep_by_slide: dict[str, dict], guardrail: GuardrailSchema
):
    if inference_service.llm is None:
        for slide in slides_data:
            slide_idx = slide.get("index", 0)
            deep_by_slide.setdefault(str(slide_idx), {})
            deep_by_slide[str(slide_idx)]["review"] = {
                "llm_understanding": "Review summary unavailable for this slide.",
                "layout_intelligence": "Layout assessment unavailable.",
                "guardrail_alignment": {
                    "status": "partial",
                    "notes": "Unable to compute alignment summary.",
                },
                "score_rationale": "Detailed rationale unavailable.",
                "detailed_recommendations": [],
                "debug_reason": "LLM unavailable for slide review generation.",
            }
        return deep_by_slide

    guardrail_context = {
        "engagement_type": guardrail.engagement_type,
        "pass_threshold": guardrail.pass_threshold,
        "language_rules": guardrail.language_rules,
        "playbook_rules": guardrail.playbook_rules[:8],
        "human_confirmed_rules": guardrail.human_confirmed_rules[:8],
    }

    async def _summarize_slide(slide: dict):
        slide_idx = slide.get("index", 0)
        deep = deep_by_slide.get(str(slide_idx), {})
        findings = []
        for agent in deep.get("agents", []):
            findings.extend(agent.get("findings", [])[:3])
        findings.extend(deep.get("judge", {}).get("findings", [])[:3])

        prompt = f"""Review this consulting slide and return a short structured JSON summary.

SLIDE TITLE: {slide.get("title", "")}
SLIDE TEXT:
{(slide.get("full_text", "") or "")[:1800]}

GUARDRAIL:
{json.dumps(guardrail_context)[:1800]}

CURRENT FINDINGS:
{json.dumps(findings)[:1800]}

Return ONLY JSON:
{{
  "llm_understanding": "<2 sentence summary of what this slide is trying to say>",
  "layout_intelligence": "<short view on layout/framework quality like a consultant reviewer>",
  "guardrail_alignment": {{
    "status": "aligned" | "partial" | "misaligned",
    "notes": "<short explanation>"
  }},
  "score_rationale": "<why this slide scored the way it did>",
  "detailed_recommendations": ["<action 1>", "<action 2>", "<action 3>"]
}}"""

        try:
            messages = [Message(role="user", content=prompt)]
            response = await inference_service.llm.generate(
                messages,
                max_tokens=min(
                    500, getattr(inference_service, "analysis_max_tokens", 500)
                ),
            )
            payload = parse_json_response(response.content)
            if isinstance(payload, dict):
                payload.setdefault("debug_reason", "")
                deep_by_slide[str(slide_idx)]["review"] = payload
            else:
                deep_by_slide[str(slide_idx)]["review"] = {
                    "llm_understanding": "Review summary unavailable for this slide.",
                    "layout_intelligence": "Layout assessment unavailable.",
                    "guardrail_alignment": {
                        "status": "partial",
                        "notes": "Unable to compute alignment summary.",
                    },
                    "score_rationale": "Detailed rationale unavailable.",
                    "detailed_recommendations": [],
                    "debug_reason": "Slide review parser returned a non-dictionary payload.",
                }
        except Exception:
            import traceback

            deep_by_slide[str(slide_idx)]["review"] = {
                "llm_understanding": "Review summary unavailable for this slide.",
                "layout_intelligence": "Layout assessment unavailable.",
                "guardrail_alignment": {
                    "status": "partial",
                    "notes": "Unable to compute alignment summary.",
                },
                "score_rationale": "Detailed rationale unavailable.",
                "detailed_recommendations": [],
                "debug_reason": traceback.format_exc().strip().splitlines()[-1][:300],
            }

    await asyncio.gather(*[_summarize_slide(slide) for slide in slides_data])
    return deep_by_slide


async def _refresh_single_slide_deep_analysis(session: dict, slide_index: int) -> dict:
    slides_data = session.get("slides_data", [])
    if slide_index < 0 or slide_index >= len(slides_data):
        raise HTTPException(status_code=404, detail="Slide not found")

    guardrail = session.get("guardrail")
    if not guardrail:
        guardrail = guardrail_manager.create_guardrail()
        session["guardrail"] = guardrail
    elif isinstance(guardrail, dict):
        guardrail = GuardrailSchema(**guardrail)
        session["guardrail"] = guardrail

    llm_settings = session.get("llm_settings", {})
    if llm_settings.get("context_window"):
        inference_service.set_local_context_window(llm_settings["context_window"])

    annotations_by_slide = session.get("annotations_by_slide", {})
    if not annotations_by_slide and session.get("scorecard"):
        annotations_by_slide = _group_annotations_by_slide(
            session.get("scorecard", {}).get("annotations", [])
        )
        session["annotations_by_slide"] = annotations_by_slide

    existing_deep = session.get("deep_analysis_by_slide", {}) or {}
    slide_key = str(slide_index)
    slide_payload = dict(
        existing_deep.get(
            slide_key,
            {
                "agents": [],
                "judge": {
                    "name": "Language Analysis",
                    "findings": [],
                    "summary": "No deep analysis available.",
                },
            },
        )
    )

    current_judge_findings = slide_payload.get("judge", {}).get("findings", [])
    non_language_findings = [
        f for f in current_judge_findings if f.get("category") == "claim_grounding"
    ]
    try:
        language_result = await language_agent.analyze_slide(
            slide_index,
            (slides_data[slide_index].get("full_text", "") or "")
            + " "
            + (slides_data[slide_index].get("title", "") or ""),
            guardrail.language_rules,
        )
        judge_findings = [
            _serialize_annotation(a)
            for a in (
                language_result.grammar_issues
                + language_result.quality_issues
                + language_result.tone_issues
            )
        ] + non_language_findings
    except Exception:
        judge_findings = current_judge_findings

    slide_payload["judge"] = {
        "name": "Language Analysis",
        "findings": judge_findings,
        "summary": (
            "Language and tone checks passed."
            if not judge_findings
            else f"{len(judge_findings)} language or tone issue(s) detected."
        ),
    }

    refreshed = await _enrich_deep_analysis_with_slide_reviews(
        [slides_data[slide_index]],
        {slide_key: slide_payload},
        guardrail,
    )
    updated_slide_payload = refreshed.get(slide_key, slide_payload)

    all_deep = dict(existing_deep)
    all_deep[slide_key] = updated_slide_payload
    session["deep_analysis_by_slide"] = all_deep

    document_fingerprint = session.get("document_fingerprint")
    deck_path = session.get("deck_path")
    if document_fingerprint and deck_path and session.get("scorecard"):
        analysis_history_store.save_analysis(
            fingerprint=document_fingerprint,
            original_filename=session.get("original_filename") or Path(deck_path).name,
            deck_path=deck_path,
            session_id=session.get("session_id") or "",
            slides_data=session.get("slides_data", []),
            scorecard=session.get("scorecard", {}),
            agent_metadata=session.get("agent_metadata", {}),
            deep_analysis_by_slide=all_deep,
            deck_metadata=session.get("deck_metadata", {}),
        )

    return updated_slide_payload


def _element_box_to_percent(element: dict, slide: dict) -> dict:
    unit = (element.get("coord_unit") or "percent").lower()
    x = float(element.get("x", 0) or 0)
    y = float(element.get("y", 0) or 0)
    width = float(element.get("width", 0) or 0)
    height = float(element.get("height", 0) or 0)

    if unit == "absolute":
        slide_w = float(slide.get("width", 0) or 0)
        slide_h = float(slide.get("height", 0) or 0)
        return {
            "top": _clamp_percent((y / slide_h) * 100 if slide_h > 0 else 0),
            "left": _clamp_percent((x / slide_w) * 100 if slide_w > 0 else 0),
            "width": _clamp_percent((width / slide_w) * 100 if slide_w > 0 else 0),
            "height": _clamp_percent((height / slide_h) * 100 if slide_h > 0 else 0),
        }

    return {
        "top": _clamp_percent(y),
        "left": _clamp_percent(x),
        "width": _clamp_percent(width),
        "height": _clamp_percent(height),
    }


def _annotation_to_bounding_box(annotation: dict, slide: dict) -> dict:
    """
    Map an annotation to a real bounding box from the slide's elements.
    Searches text_boxes, charts, tables, and images by matching text or shape_id.
    Returns percentage-based coordinates for overlay rendering.
    """
    ann_text = (annotation.get("text") or "").lower().strip()
    ann_shape_id = annotation.get("shape_id")
    ann_category = annotation.get("category", "")

    slide_w = slide.get("width", 10.0)
    slide_h = slide.get("height", 7.5)

    def _to_percent_box(tb: dict) -> dict:
        unit = (tb.get("coord_unit") or "absolute").lower()
        x_val = float(tb.get("x", 0) or 0)
        y_val = float(tb.get("y", 0) or 0)
        w_val = float(tb.get("width", 0) or 0)
        h_val = float(tb.get("height", 0) or 0)
        if unit == "percent":
            return {
                "top": _clamp_percent(y_val),
                "left": _clamp_percent(x_val),
                "width": _clamp_percent(w_val),
                "height": _clamp_percent(h_val),
            }
        return {
            "top": _clamp_percent((y_val / slide_h) * 100 if slide_h > 0 else 0),
            "left": _clamp_percent((x_val / slide_w) * 100 if slide_w > 0 else 0),
            "width": _clamp_percent((w_val / slide_w) * 100 if slide_w > 0 else 50),
            "height": _clamp_percent((h_val / slide_h) * 100 if slide_h > 0 else 10),
        }

    candidate_text_boxes = list(slide.get("text_boxes", []) or [])
    title_text = (slide.get("title") or "").strip()
    if title_text:
        title_exists = any(
            (tb.get("text") or "").strip().lower() == title_text.lower()
            for tb in candidate_text_boxes
        )
        if not title_exists:
            candidate_text_boxes.insert(
                0,
                {
                    "id": "title_proxy",
                    "text": title_text,
                    "x": 4.0,
                    "y": 3.0,
                    "width": 92.0,
                    "height": 12.0,
                    "coord_unit": "percent",
                },
            )

    # 1. Try to match by shape_id (PPTX)
    if ann_shape_id:
        for tb in candidate_text_boxes:
            if tb.get("id") == ann_shape_id or ann_shape_id in tb.get("id", ""):
                return _to_percent_box(tb)

    # 2. Try to match by text content in text boxes
    if ann_text:
        for tb in candidate_text_boxes:
            tb_text = (tb.get("text") or "").lower()
            if ann_text[:30] in tb_text or tb_text[:30] in ann_text:
                return _to_percent_box(tb)

    # 3. If data_accuracy/chart category, match to chart elements
    if "data" in ann_category or "chart" in ann_text:
        for chart in slide.get("charts", []):
            chart_id = chart.get("id")
            chart_title = (chart.get("title") or "").lower()
            if (
                ann_shape_id
                and chart_id
                and (ann_shape_id == chart_id or ann_shape_id in str(chart_id))
            ):
                return _element_box_to_percent(chart, slide)
            if ann_text and (ann_text[:20] in chart_title or chart_title in ann_text):
                if chart.get("width", 0) > 0 and chart.get("height", 0) > 0:
                    return _element_box_to_percent(chart, slide)
                return {"top": 30, "left": 15, "width": 70, "height": 50}

    # 4. If visual/table category, match to table elements
    if "table" in ann_text or "table" in ann_category:
        for tbl in slide.get("tables", []):
            tbl_id = tbl.get("id")
            if (
                ann_shape_id
                and tbl_id
                and (ann_shape_id == tbl_id or ann_shape_id in str(tbl_id))
            ):
                return _element_box_to_percent(tbl, slide)
            if tbl.get("x", 0) > 0:
                return _element_box_to_percent(tbl, slide)

    # 5. If image-related
    if "image" in ann_text or "figure" in ann_text:
        for img in slide.get("images", []):
            img_id = img.get("id")
            if (
                ann_shape_id
                and img_id
                and (ann_shape_id == img_id or ann_shape_id in str(img_id))
            ):
                return _element_box_to_percent(img, slide)
            if img.get("x", 0) > 0:
                return _element_box_to_percent(img, slide)

    # 6. Fallback: distribute fixes vertically to avoid overlap
    fix_index = _deterministic_band_index(ann_text, modulo=6)
    return {
        "top": 4 + fix_index * 15,
        "left": 2,
        "width": 96,
        "height": 10,
    }


# Ensure static directory exists
static_dir = Path(__file__).parent.parent / "static"
static_dir.mkdir(parents=True, exist_ok=True)

# Mount /assets to serve CSS/JS directly
assets_dir = static_dir / "assets"
if not assets_dir.exists():
    assets_dir.mkdir(parents=True, exist_ok=True)
app.mount("/assets", StaticFiles(directory=str(assets_dir)), name="assets")


async def health_check():
    ensure_services_loaded()

    lt_status = await language_agent.lt_client.status()
    grammar_available = bool(lt_status.get("available", False))
    llm_available = inference_service.llm is not None
    vision_backend = getattr(
        getattr(analysis_orchestrator, "visual_agent", None),
        "vision_backend",
        "unknown",
    )
    chroma_ok = chroma_manager is not None and chroma_manager.client is not None
    surya_layout_ready = model_registry.get_surya_layout() is not None
    surya_recognition_ready = model_registry.get_surya_recognition() is not None
    degraded = []
    if not llm_available:
        degraded.append("llm_unavailable")
    if not grammar_available:
        degraded.append("grammar_regex_fallback")
    if vision_backend == "fallback":
        degraded.append("vision_fallback")
    if not surya_layout_ready:
        degraded.append("layout_fallback")
    if not surya_recognition_ready:
        degraded.append("ocr_fallback")
    if not chroma_ok:
        degraded.append("evidence_store_unavailable")

    return {
        "status": "healthy" if not degraded else "degraded",
        "offline_mode": True,
        "components": {
            "llm": llm_available,
            "grammar": {
                "available": grammar_available,
                "base_url": language_agent.lt_client.base_url,
                "last_error": lt_status.get("last_error"),
            },
            "vision": vision_backend,
            "surya": {
                "layout": surya_layout_ready,
                "recognition": surya_recognition_ready,
                "last_error": model_registry._surya_last_error,
            },
            "chroma": chroma_ok,
        },
        "degraded_reasons": degraded,
    }


async def get_session_metrics(
    include_sessions: bool = False, x_user_role: str = Header("junior")
):
    if x_user_role != "senior":
        raise HTTPException(
            status_code=403, detail="Senior evaluator role required for this action"
        )

    now_ts = time.time()
    ttl_seconds = _session_ttl_seconds()
    sessions_snapshot = list(active_sessions.items())

    sessions_by_namespace = _count_sessions_by_namespace()
    expiring_within_1h = 0
    idle_values = []
    age_values = []

    session_rows = []
    for sid, session in sessions_snapshot:
        last_access_ts = _session_last_access_ts(session, now_ts)
        created_at_ts = float(session.get("created_at_ts", last_access_ts))
        idle_seconds = max(0, int(now_ts - last_access_ts))
        age_seconds = max(0, int(now_ts - created_at_ts))
        ttl_remaining_seconds = max(0, int(ttl_seconds - idle_seconds))
        if 0 < ttl_remaining_seconds <= 3600:
            expiring_within_1h += 1

        idle_values.append(idle_seconds)
        age_values.append(age_seconds)

        if include_sessions:
            session_rows.append(
                {
                    "session_id": sid,
                    "status": session.get("status", "unknown"),
                    "client_namespace": session.get("client_namespace"),
                    "namespace_key": session.get("namespace_key")
                    or _namespace_key(session.get("client_namespace")),
                    "idle_seconds": idle_seconds,
                    "age_seconds": age_seconds,
                    "ttl_remaining_seconds": ttl_remaining_seconds,
                    "has_deck": bool(session.get("deck_path")),
                    "slide_count": len(session.get("slides_data", [])),
                }
            )

    with _cleanup_stats_lock:
        cleanup_stats = dict(_cleanup_stats)

    payload = {
        "timestamp": utc_now_iso(),
        "limits": {
            "session_ttl_hours": settings.session_ttl_hours,
            "session_cleanup_interval_minutes": settings.session_cleanup_interval_minutes,
            "max_active_sessions": settings.max_active_sessions,
            "max_active_sessions_per_namespace": settings.max_active_sessions_per_namespace,
        },
        "sessions": {
            "active_count": len(sessions_snapshot),
            "expiring_within_1h_count": expiring_within_1h,
            "by_namespace": sessions_by_namespace,
            "max_idle_seconds": max(idle_values) if idle_values else 0,
            "max_age_seconds": max(age_values) if age_values else 0,
        },
        "cleanup": {
            "last_run_ts": cleanup_stats.get("last_run_ts"),
            "last_expired_count": cleanup_stats.get("last_expired_count", 0),
            "last_expired_ids": cleanup_stats.get("last_expired_ids", []),
            "total_expired_count": cleanup_stats.get("total_expired_count", 0),
        },
    }

    if include_sessions:
        payload["session_rows"] = session_rows

    return payload


async def get_recent_analysis_history(limit: int = 12):
    ensure_services_loaded()
    limit = max(1, min(limit, 50))
    return {"items": analysis_history_store.list_recent(limit)}


async def open_analysis_history(fingerprint: str):
    ensure_services_loaded()
    history_entry = analysis_history_store.get_by_fingerprint(fingerprint)
    if history_entry is None:
        raise HTTPException(status_code=404, detail="History entry not found")
    if not _history_entry_is_current(history_entry):
        raise HTTPException(
            status_code=409,
            detail="Saved analysis is outdated. Re-upload the document to regenerate with the latest analysis logic.",
        )

    import uuid

    session_id = str(uuid.uuid4())
    now_ts = time.time()
    active_sessions[session_id] = {
        "client_namespace": None,
        "namespace_key": _namespace_key(None),
        "deck_path": None,
        "slides_data": [],
        "source_namespace": None,
        "source_files": [],
        "source_indexed_chunks": {},
        "excel_data": None,
        "guardrail": None,
        "scorecard": None,
        "status": "created",
        "created_at_ts": now_ts,
        "last_access_ts": now_ts,
    }
    session = active_sessions[session_id]
    await _restore_history_to_session(session_id, session, history_entry)
    _persist_session(session_id, session)
    return {
        "session_id": session_id,
        "slide_count": len(session.get("slides_data", [])),
        "status": session.get("status"),
        "restored_from_history": True,
    }


@app.post("/api/session/create")
async def create_session(client_namespace: Optional[str] = None):
    import uuid

    session_id = str(uuid.uuid4())
    now_ts = time.time()
    namespace_key = _namespace_key(client_namespace)
    _enforce_session_capacity(namespace_key)

    active_sessions[session_id] = {
        "client_namespace": client_namespace,
        "namespace_key": namespace_key,
        "deck_path": None,
        "slides_data": [],
        "source_namespace": None,
        "source_files": [],
        "source_indexed_chunks": {},
        "excel_data": None,
        "guardrail": None,
        "scorecard": None,
        "status": "created",
        "created_at_ts": now_ts,
        "last_access_ts": now_ts,
        "llm_settings": {
            "provider": inference_service.current_provider,
            "context_window": inference_service.local_context_window,
        },
    }

    _persist_session(session_id, active_sessions[session_id])

    return {"session_id": session_id}


def require_senior(x_user_role: str = Header("junior")):
    if x_user_role != "senior":
        raise HTTPException(
            status_code=403, detail="Senior evaluator role required for this action"
        )
    return x_user_role


def _slugify_guardrail_template_name(name: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", (name or "").strip()).strip("_").lower()
    if not cleaned:
        raise HTTPException(status_code=400, detail="Template name cannot be empty")
    return cleaned[:80]


def _guardrail_template_dir() -> Path:
    template_dir = data_dir / "guardrail_templates"
    template_dir.mkdir(parents=True, exist_ok=True)
    return template_dir


def _resolve_guardrail_template_file(file_id: str) -> Path:
    template_dir = _guardrail_template_dir().resolve()
    candidate = (template_dir / Path(file_id).name).resolve()
    if template_dir not in candidate.parents or not candidate.exists():
        raise HTTPException(
            status_code=400, detail="Invalid guardrail template reference"
        )
    return candidate


def _apply_guardrail_to_session_state(
    session: dict, guardrail: GuardrailSchema
) -> tuple[GuardrailSchema, bool]:
    existing_guardrail = session.get("guardrail")
    applied_guardrail = GuardrailSchema(**guardrail.model_dump())
    applied_guardrail.client_namespace = (
        applied_guardrail.client_namespace or session.get("client_namespace")
    )
    applied_guardrail.engagement_type = applied_guardrail.engagement_type or (
        existing_guardrail.engagement_type
        if isinstance(existing_guardrail, GuardrailSchema)
        else "strategy"
    )

    signature_cleared = bool(
        applied_guardrail.signed_by
        or applied_guardrail.signed_at
        or applied_guardrail.sha256
        or applied_guardrail.signature
        or applied_guardrail.public_key
        or (
            isinstance(existing_guardrail, GuardrailSchema)
            and (
                existing_guardrail.signed_by
                or existing_guardrail.signed_at
                or existing_guardrail.sha256
                or existing_guardrail.signature
                or existing_guardrail.public_key
            )
        )
    )
    applied_guardrail.signed_by = None
    applied_guardrail.signed_at = None
    applied_guardrail.sha256 = None
    applied_guardrail.signature = None
    applied_guardrail.public_key = None
    applied_guardrail.signature_algorithm = None

    session["guardrail"] = applied_guardrail
    session["guardrail_updated_at"] = utc_now_iso()

    discovery_state = session.get("discovery_state")
    if isinstance(discovery_state, dict):
        discovery_state["applied_at"] = session["guardrail_updated_at"]

    _persist_known_session(session)

    return applied_guardrail, signature_cleared


def _guardrail_from_session(session: dict) -> GuardrailSchema:
    raw_guardrail = session.get("guardrail")
    if isinstance(raw_guardrail, GuardrailSchema):
        return raw_guardrail
    if isinstance(raw_guardrail, dict):
        try:
            return GuardrailSchema(**raw_guardrail)
        except Exception:
            pass
    return GuardrailSchema(client_namespace=session.get("client_namespace"))


def _stringify_guardrail_rule(rule: object) -> str:
    if isinstance(rule, str):
        return rule.strip()
    if isinstance(rule, dict):
        for key in ("rule", "name", "description", "text", "pattern", "title"):
            value = rule.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
        try:
            return json.dumps(rule)[:180]
        except Exception:
            return str(rule)[:180]
    return str(rule)[:180]


def _rule_categories_from_text(rule_text: str) -> set[str]:
    text = (rule_text or "").lower()
    categories: set[str] = set()

    if any(token in text for token in ("claim", "source", "citation", "evidence")):
        categories.update({"claim_grounding", "claim_extraction"})
    if any(token in text for token in ("data", "excel", "lineage", "metric", "number")):
        categories.update({"data_accuracy", "benchmarking"})
    if any(
        token in text
        for token in ("grammar", "tone", "style", "language", "wording", "headline")
    ):
        categories.update({"language", "grammar", "tone", "quality", "hedging"})
    if any(token in text for token in ("visual", "chart", "table", "layout", "design")):
        categories.update({"visual", "structure"})
    if any(
        token in text
        for token in ("framework", "so what", "recommendation", "storyline")
    ):
        categories.update({"framework", "so_what", "structure"})

    if not categories:
        categories.update(
            {
                "claim_grounding",
                "data_accuracy",
                "language",
                "visual",
                "structure",
                "framework",
                "so_what",
                "benchmarking",
            }
        )

    return categories


def _build_guardrail_coverage(
    session: dict,
    slide: dict,
    slide_annotations: list[dict],
    slide_score: int,
) -> list[dict]:
    guardrail = _guardrail_from_session(session)
    coverage: list[dict] = []

    failing_annotations = [
        annotation
        for annotation in slide_annotations
        if annotation.get("severity") in ("hard_block", "warning")
    ]

    def add_item(
        item_id: str,
        rule: str,
        source: str,
        status: str,
        detail: str,
    ) -> None:
        coverage.append(
            {
                "id": item_id,
                "rule": rule,
                "source": source,
                "status": status,
                "detail": detail,
            }
        )

    threshold = int(guardrail.pass_threshold or 75)
    threshold_status = "checked" if slide_score >= threshold else "failed"
    add_item(
        "pass-threshold",
        f"Pass threshold >= {threshold}",
        "system",
        threshold_status,
        (
            f"Slide score {slide_score} meets the configured threshold."
            if threshold_status == "checked"
            else f"Slide score {slide_score} is below threshold {threshold}."
        ),
    )

    source_files = session.get("source_files") or []
    claim_failures = [
        annotation
        for annotation in failing_annotations
        if annotation.get("category") in {"claim_grounding", "claim_extraction"}
    ]
    if source_files:
        add_item(
            "source-grounding",
            "Claim grounding against uploaded sources",
            "system",
            "failed" if claim_failures else "checked",
            (
                claim_failures[0].get("message", "Claim grounding issue detected.")
                if claim_failures
                else "No claim-grounding violations detected for this slide."
            ),
        )
    else:
        add_item(
            "source-grounding",
            "Claim grounding against uploaded sources",
            "system",
            "skipped",
            "No source documents uploaded for this session.",
        )

    excel_data = session.get("excel_data")
    lineage_failures = [
        annotation
        for annotation in failing_annotations
        if annotation.get("category") in {"data_accuracy", "benchmarking"}
    ]
    if excel_data:
        add_item(
            "excel-lineage",
            "Data lineage against uploaded Excel source",
            "system",
            "failed" if lineage_failures else "checked",
            (
                lineage_failures[0].get("message", "Data lineage issue detected.")
                if lineage_failures
                else "No data-lineage violations detected for this slide."
            ),
        )
    else:
        add_item(
            "excel-lineage",
            "Data lineage against uploaded Excel source",
            "system",
            "skipped",
            "No Excel source uploaded for this session.",
        )

    for idx, rule in enumerate((guardrail.human_confirmed_rules or [])[:8]):
        rule_text = _stringify_guardrail_rule(rule)
        categories = _rule_categories_from_text(rule_text)
        matched_failures = [
            annotation
            for annotation in failing_annotations
            if annotation.get("category") in categories
        ]
        add_item(
            f"human-{idx}",
            rule_text or f"Human rule {idx + 1}",
            "human_confirmed",
            "failed" if matched_failures else "checked",
            (
                matched_failures[0].get("message", "Violation detected for this rule.")
                if matched_failures
                else "No violations detected for this rule on this slide."
            ),
        )

    for idx, rule in enumerate((guardrail.playbook_rules or [])[:8]):
        rule_text = _stringify_guardrail_rule(rule)
        categories = _rule_categories_from_text(rule_text)
        matched_failures = [
            annotation
            for annotation in failing_annotations
            if annotation.get("category") in categories
        ]
        add_item(
            f"playbook-{idx}",
            rule_text or f"Playbook rule {idx + 1}",
            "playbook",
            "failed" if matched_failures else "checked",
            (
                matched_failures[0].get("message", "Violation detected for this rule.")
                if matched_failures
                else "No violations detected for this rule on this slide."
            ),
        )

    language_rules = guardrail.language_rules or {}
    prohibited_phrases = language_rules.get("prohibited_phrases", [])
    if isinstance(prohibited_phrases, list) and prohibited_phrases:
        combined_text = f"{slide.get('title', '')} {slide.get('full_text', '')}".lower()
        matched_phrases = [
            phrase
            for phrase in prohibited_phrases
            if isinstance(phrase, str)
            and phrase.strip()
            and phrase.lower() in combined_text
        ]
        add_item(
            "language-prohibited-phrases",
            "Language prohibited phrases",
            "language",
            "failed" if matched_phrases else "checked",
            (
                f"Detected prohibited phrases: {', '.join(matched_phrases[:3])}"
                if matched_phrases
                else "No prohibited phrases detected on this slide."
            ),
        )

    language_rule_keys = [
        key for key in language_rules.keys() if key != "prohibited_phrases"
    ]
    for idx, key in enumerate(language_rule_keys[:6]):
        add_item(
            f"language-{idx}",
            f"Language rule: {key}",
            "language",
            "skipped",
            "Rule is active, but deterministic per-slide validation is not yet available.",
        )

    visual_patterns = (guardrail.discovered_patterns or {}).get("visual")
    if isinstance(visual_patterns, dict) and visual_patterns:
        visual_failures = [
            annotation
            for annotation in failing_annotations
            if annotation.get("category") in {"visual", "structure"}
        ]
        add_item(
            "discovered-visual-patterns",
            "Discovered visual pattern checks",
            "discovered",
            "failed" if visual_failures else "checked",
            (
                visual_failures[0].get("message", "Visual rule issue detected.")
                if visual_failures
                else "No discovered visual-pattern violations detected on this slide."
            ),
        )

    return coverage


@app.get("/api/session/{session_id}/evidence")
async def get_session_evidence(session_id: str):
    session = _get_session_or_404(session_id)
    source_files = session.get("source_files") or []
    indexed_chunks = session.get("source_indexed_chunks") or {}
    evidence_sources = [
        {
            "filename": filename,
            "documents_indexed": int(indexed_chunks.get(filename, 0)),
        }
        for filename in source_files
    ]

    excel_data = session.get("excel_data") or {}
    sheets = excel_data.get("sheets") or {}
    excel_snapshot = None
    if sheets:
        file_path = excel_data.get("file_path")
        excel_snapshot = {
            "filename": Path(file_path).name if file_path else "uploaded.xlsx",
            "sheets": list(sheets.keys()),
            "total_rows": sum(len(rows) for rows in sheets.values()),
        }

    return {
        "session_id": session_id,
        "evidence_sources": evidence_sources,
        "excel_snapshot": excel_snapshot,
        "source_namespace": session.get("source_namespace"),
    }


@app.post("/api/session/{session_id}/upload")
async def upload_deck(
    session_id: str, background_tasks: BackgroundTasks, file: UploadFile = File(...)
):
    session = _get_session_or_404(session_id)
    ensure_services_loaded()

    # MIME Validation & Extension Hardening (Slides only)
    allowed_slide_exts = {".pdf", ".pptx"}
    filename = Path(file.filename or "deck.pdf").name
    ext = os.path.splitext(filename)[1].lower()

    if ext not in allowed_slide_exts:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file extension: {ext}. For analysis, only PDF and PPTX are allowed.",
        )

    # Basic MIME check (heuristic)
    content_type = (file.content_type or "").strip()
    logger.info(
        f"Upload MIME type: '{content_type}' (repr: {repr(content_type)}), filename: {filename}, ext: {ext}"
    )

    # Explicit allowlist of known valid MIME types for slides
    allowed_mime_types = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "application/octet-stream",
        "",  # empty content_type is OK, we rely on extension check above
    }

    # Accept if MIME type is in our allowlist, starts with "application/", or is empty
    if (
        content_type
        and content_type not in allowed_mime_types
        and not content_type.startswith("application/")
    ):
        raise HTTPException(
            status_code=400, detail=f"Invalid MIME type: {content_type}"
        )

    upload_dir = data_dir / "uploads" / session_id

    upload_dir.mkdir(parents=True, exist_ok=True)

    file_path = upload_dir / filename

    total_written = 0
    chunk_size = 1024 * 1024
    max_size = int(settings.max_file_size)
    with open(file_path, "wb") as buffer:
        while True:
            chunk = await file.read(chunk_size)
            if not chunk:
                break
            total_written += len(chunk)
            if total_written > max_size:
                buffer.close()
                try:
                    file_path.unlink(missing_ok=True)
                except Exception:
                    pass
                raise HTTPException(
                    status_code=413,
                    detail=f"Uploaded file exceeds max size of {max_size} bytes",
                )
            buffer.write(chunk)

    document_fingerprint = analysis_history_store.compute_sha256(str(file_path))
    history_entry = analysis_history_store.get_by_fingerprint(document_fingerprint)

    session["deck_path"] = str(file_path)
    session["document_fingerprint"] = document_fingerprint
    session["original_filename"] = filename
    session["history_available"] = history_entry is not None
    session["status"] = "uploaded"
    _persist_session(session_id, session)

    return {
        "filename": filename,
        "path": str(file_path),
        "session_id": session_id,
        "document_fingerprint": document_fingerprint,
        "history_available": history_entry is not None,
    }


@app.post("/api/session/{session_id}/analyze")
async def analyze_deck(session_id: str):
    session = _get_session_or_404(session_id)
    deck_path = session.get("deck_path")

    if not deck_path or not os.path.exists(deck_path):
        raise HTTPException(status_code=400, detail="No deck uploaded")

    try:
        ensure_services_loaded()
        llm_settings = session.get("llm_settings", {})
        if llm_settings.get("context_window"):
            inference_service.set_local_context_window(llm_settings["context_window"])
        document_fingerprint = session.get("document_fingerprint")
        history_entry = (
            analysis_history_store.get_by_fingerprint(document_fingerprint)
            if document_fingerprint
            else None
        )
        upload_dir = Path(deck_path).parent
        previews_dir = upload_dir / "previews"
        previews_dir.mkdir(parents=True, exist_ok=True)
        coord_unit = "absolute" if deck_path.lower().endswith(".pptx") else "percent"

        if _history_entry_is_current(history_entry) and history_entry.get("scorecard"):
            await _generate_previews_for_deck(deck_path, upload_dir)
            session.update(
                {
                    "deck_metadata": history_entry.get("deck_metadata", {}),
                    "slides_data": _clone_slide_payloads_for_session(
                        session_id, history_entry.get("slides_data", []), upload_dir
                    ),
                    "scorecard": history_entry.get("scorecard", {}),
                    "annotations_by_slide": _group_annotations_by_slide(
                        history_entry.get("scorecard", {}).get("annotations", [])
                    ),
                    "agent_metadata": history_entry.get("agent_metadata", {}),
                    "deep_analysis_by_slide": history_entry.get(
                        "deep_analysis_by_slide", {}
                    ),
                    "status": "analyzed",
                    "history_restored": True,
                    "history_restored_at": utc_now_iso(),
                }
            )
            _persist_session(session_id, session)
            return {
                "session_id": session_id,
                "slide_count": len(session.get("slides_data", [])),
                "status": "parsed",
                "restored_from_history": True,
            }

        if deck_path.lower().endswith(".pptx"):
            deck_content = await ingestion_service.ingest_pptx(deck_path)
            await ingestion_service.convert_pptx_to_images(deck_path, str(previews_dir))
        elif deck_path.lower().endswith(".pdf"):
            deck_content = await ingestion_service.ingest_pdf(deck_path)
            await ingestion_service.convert_pdf_to_images(deck_path, str(previews_dir))
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format")

        slides_data = []
        for slide in deck_content.slides:
            preview_path = str(previews_dir / f"slide_{slide.slide_index}.png")
            slide_assets_dir = upload_dir / "assets"
            slide_assets_dir.mkdir(parents=True, exist_ok=True)
            image_payloads = []
            for img in slide.images:
                asset_path = None
                asset_url = None
                extension = (img.extension or "bin").lower().lstrip(".")
                content_type = img.content_type or "application/octet-stream"
                if img.image_data:
                    asset_filename = f"slide_{slide.slide_index}_{img.id}.{extension}"
                    asset_file = slide_assets_dir / asset_filename
                    asset_file.write_bytes(img.image_data)
                    asset_path = str(asset_file)
                    asset_url = f"/api/session/{session_id}/slide/{slide.slide_index}/asset/{img.id}"

                image_payloads.append(
                    {
                        "id": img.id,
                        "x": img.x,
                        "y": img.y,
                        "width": img.width,
                        "height": img.height,
                        "has_content": img.image_data is not None,
                        "coord_unit": coord_unit,
                        "asset_path": asset_path,
                        "asset_url": asset_url,
                        "content_type": content_type,
                        "extension": extension if img.image_data else None,
                    }
                )
            slides_data.append(
                {
                    "id": f"slide_{slide.slide_index}",
                    "index": slide.slide_index,
                    "title": slide.title,
                    "preview_path": preview_path,
                    "previewUrl": f"/api/session/{session_id}/slide/{slide.slide_index}/image",
                    "full_text": slide.title
                    + " "
                    + " ".join(tb.text for tb in slide.text_boxes),
                    "text_boxes": [
                        {
                            "id": tb.id,
                            "text": tb.text,
                            "x": tb.x,
                            "y": tb.y,
                            "width": tb.width,
                            "height": tb.height,
                            "runs": [
                                {
                                    "text": r.text,
                                    "font_name": r.font_name,
                                    "font_size": r.font_size,
                                    "font_bold": r.font_bold,
                                }
                                for r in tb.runs
                            ],
                        }
                        for tb in slide.text_boxes
                    ],
                    "charts": [
                        {
                            "id": c.chart_id,
                            "title": c.title,
                            "type": c.chart_type,
                            "x": c.x,
                            "y": c.y,
                            "width": c.width,
                            "height": c.height,
                            "coord_unit": coord_unit,
                            "data_range_ref": c.data_range_ref,
                            "cache_values": c.cache_values,
                        }
                        for c in slide.charts
                    ],
                    "tables": [
                        {
                            "id": t.table_id,
                            "rows": t.rows,
                            "columns": t.columns,
                            "title": t.title,
                            "text": t.text,
                            "x": t.x,
                            "y": t.y,
                            "width": t.width,
                            "height": t.height,
                            "coord_unit": coord_unit,
                        }
                        for t in slide.tables
                    ],
                    "images": image_payloads,
                    "width": slide.width,
                    "height": slide.height,
                    "ocr_backend": slide.ocr_backend,
                }
            )

        session["slides_data"] = slides_data
        session["deck_metadata"] = deck_content.metadata
        session["status"] = "parsed"
        _persist_session(session_id, session)

        audit_log_service.log_event(
            session_id=session_id,
            user_role="system",
            action="PARSE_COMPLETE",
            details={
                "slide_count": len(slides_data),
                "format": "PDF" if deck_path.endswith(".pdf") else "PPTX",
            },
        )

        return {
            "session_id": session_id,
            "slide_count": len(slides_data),
            "status": "parsed",
        }
    except Exception as e:
        import traceback

        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/session/{session_id}/slide/{index}/image")
async def get_slide_image(session_id: str, index: int):
    session = _get_session_or_404(session_id)
    deck_path = session.get("deck_path")
    if not deck_path:
        raise HTTPException(status_code=404, detail="No deck found for session")

    upload_dir = Path(deck_path).parent
    image_path = upload_dir / "previews" / f"slide_{index}.png"

    if not image_path.exists():
        raise HTTPException(status_code=404, detail="Slide image not found")

    return FileResponse(str(image_path), media_type="image/png")


@app.get("/api/session/{session_id}/slide/{index}/asset/{image_id}")
async def get_slide_asset(session_id: str, index: int, image_id: str):
    session = _get_session_or_404(session_id)
    slides_data = session.get("slides_data", [])
    if index < 0 or index >= len(slides_data):
        raise HTTPException(status_code=404, detail="Slide not found")

    slide = slides_data[index]
    image_items = slide.get("images", []) or []
    matching = next(
        (item for item in image_items if str(item.get("id")) == image_id), None
    )
    if not matching:
        raise HTTPException(status_code=404, detail="Slide asset not found")

    asset_path = matching.get("asset_path")
    if not asset_path:
        raise HTTPException(status_code=404, detail="Slide asset file unavailable")

    candidate = Path(asset_path).resolve()
    upload_root = (data_dir / "uploads").resolve()
    if upload_root not in candidate.parents or not candidate.exists():
        raise HTTPException(status_code=404, detail="Slide asset file unavailable")

    return FileResponse(
        str(candidate),
        media_type=matching.get("content_type") or "application/octet-stream",
    )


@app.post("/api/session/{session_id}/run-analysis")
async def run_analysis(session_id: str):
    session = _get_session_or_404(session_id)
    slides_data = session.get("slides_data", [])

    if not slides_data:
        raise HTTPException(status_code=400, detail="No slides to analyze")

    if (
        session.get("history_restored")
        and session.get("scorecard")
        and session.get("status") == "analyzed"
    ):
        return {
            "session_id": session_id,
            "scorecard": session.get("scorecard"),
            "restored_from_history": True,
        }

    ensure_services_loaded()
    llm_settings = session.get("llm_settings", {})
    if llm_settings.get("context_window"):
        inference_service.set_local_context_window(llm_settings["context_window"])
    guardrail = session.get("guardrail")
    if not guardrail:
        guardrail = guardrail_manager.create_guardrail()
        session["guardrail"] = guardrail

    excel_data = session.get("excel_data")

    agent_results = await analysis_orchestrator.run_parallel_analysis(
        slides_data, guardrail, excel_data
    )

    language_annotations = await language_agent.analyze_deck(
        slides_data, guardrail.language_rules
    )
    source_grounding_annotations = await _build_source_grounding_annotations(session)
    language_annotations.extend(source_grounding_annotations)

    context_result = await analysis_orchestrator.run_slide_context_synthesizer(
        slides_data,
        agent_results,
        language_annotations,
    )
    agent_results.append(context_result)

    scorecard = qa_grader.calculate_scorecard(
        agent_results, language_annotations, guardrail
    )

    scorecard_data = scorecard.model_dump()
    scorecard_data["weights"] = _normalize_guardrail_weights(guardrail)
    session["scorecard"] = scorecard_data
    session["annotations_by_slide"] = _group_annotations_by_slide(
        scorecard_data.get("annotations", [])
    )
    deep_analysis_by_slide = _build_deep_analysis_by_slide(
        agent_results, language_annotations, slides_data
    )
    deep_analysis_by_slide = await _enrich_deep_analysis_with_slide_reviews(
        slides_data, deep_analysis_by_slide, guardrail
    )
    session["deep_analysis_by_slide"] = deep_analysis_by_slide
    session["status"] = "analyzed"

    # Store agent-specific metadata for detailed slide exploration
    agent_metadata = {}
    for res in agent_results:
        agent_metadata[res.agent_name] = res.metadata
    session["agent_metadata"] = agent_metadata
    _persist_session(session_id, session)

    # Cleanup VRAM after heavy deck analysis
    inference_service.optimize_memory()

    document_fingerprint = session.get("document_fingerprint")
    deck_path = session.get("deck_path")
    if document_fingerprint and deck_path:
        analysis_history_store.save_analysis(
            fingerprint=document_fingerprint,
            original_filename=session.get("original_filename") or Path(deck_path).name,
            deck_path=deck_path,
            session_id=session_id,
            slides_data=slides_data,
            scorecard=scorecard_data,
            agent_metadata=agent_metadata,
            deep_analysis_by_slide=session.get("deep_analysis_by_slide", {}),
            deck_metadata=session.get("deck_metadata", {}),
        )

    return {
        "session_id": session_id,
        "scorecard": scorecard_data,
    }


@app.get("/api/session/{session_id}/slide/{index}/analysis")
async def get_slide_analysis(session_id: str, index: int):
    session = _get_session_or_404(session_id)
    scorecard_data = session.get("scorecard")
    slides_data = session.get("slides_data", [])
    agent_metadata = session.get("agent_metadata", {})

    if index < 0:
        raise HTTPException(status_code=400, detail="Slide index must be non-negative")

    if not scorecard_data or index >= len(slides_data):
        # Allow fallback to basic data if analysis not run yet
        slide = slides_data[index] if index < len(slides_data) else {}
        fallback_score = 0
        fallback_annotations = []
        guardrail_coverage = _build_guardrail_coverage(
            session=session,
            slide=slide,
            slide_annotations=fallback_annotations,
            slide_score=fallback_score,
        )
        return {
            "id": f"slide-{index}",
            "title": slide.get("title", "Untitled"),
            "summary": "Analysis pending...",
            "overallScore": fallback_score,
            "density": "Medium",
            "visuals": [],
            "fixes": [],
            "councilDebate": [],
            "frameworkDetected": None,
            "frameworkAnalysis": None,
            "soWhatResult": None,
            "benchmarkAnalysis": None,
            "slideContext": None,
            "citationIssues": [],
            "deepAnalysis": {
                "agents": [],
                "judge": {
                    "name": "Language Analysis",
                    "findings": [],
                    "summary": "Analysis pending...",
                },
            },
            "guardrailCoverage": guardrail_coverage,
            "analysisBackends": {
                "surya": False,
                "vision": "unknown",
                "ocr": "unknown",
            },
        }

    slide = slides_data[index]
    deep_analysis_by_slide = session.get("deep_analysis_by_slide", {})

    # Filter annotations for this slide
    grouped_annotations = session.get("annotations_by_slide", {})
    if grouped_annotations:
        slide_annotations = grouped_annotations.get(str(index), [])
    else:
        all_annotations = scorecard_data.get("annotations", [])
        grouped_annotations = _group_annotations_by_slide(all_annotations)
        session["annotations_by_slide"] = grouped_annotations
        slide_annotations = grouped_annotations.get(str(index), [])

    # Get visuals and density from Visual Analysis Agent metadata (SURYA powered)
    visual_agent_meta = agent_metadata.get("Visual Analysis Agent", {})
    visual_meta = (visual_agent_meta.get("slides_analysis", {}) or {}).get(
        str(index), {}
    )
    raw_image_analysis = [dict(item) for item in visual_meta.get("image_analysis", [])]
    label_counters: dict[str, int] = {}
    keyed_image_analysis: list[dict] = []
    for item in raw_image_analysis:
        visual_key = item.get("id")
        if not visual_key:
            label = str(item.get("label") or item.get("type") or "visual").lower()
            label_counters[label] = label_counters.get(label, 0) + 1
            visual_key = f"{label}_{label_counters[label] - 1}"
        item["visualKey"] = visual_key
        keyed_image_analysis.append(item)

    surya_matchers: dict[str, list[str]] = {}
    for item in keyed_image_analysis:
        if item.get("type") == "surya_block":
            label = str(item.get("label") or "visual").lower()
            surya_matchers.setdefault(label, []).append(item["visualKey"])

    assigned_surya: dict[str, int] = {label: 0 for label in surya_matchers}
    visuals = []
    for raw_visual in visual_meta.get("visuals", []):
        visual = dict(raw_visual)
        if all(k in visual for k in ("top", "left", "width", "height")):
            visual = {
                "x": float(visual.get("left", 0) or 0),
                "y": float(visual.get("top", 0) or 0),
                "width": float(visual.get("width", 0) or 0),
                "height": float(visual.get("height", 0) or 0),
                "coord_unit": visual.get("coord_unit", "percent"),
                "label": visual.get("label", "Visual"),
                "visual_key": visual.get("visual_key"),
            }
        box = _element_box_to_percent(visual, slide)
        visual_key = visual.get("visual_key")
        if not visual_key:
            label = str(visual.get("label") or "visual").lower()
            matches = surya_matchers.get(label, [])
            next_idx = assigned_surya.get(label, 0)
            if next_idx < len(matches):
                visual_key = matches[next_idx]
                assigned_surya[label] = next_idx + 1
            else:
                visual_key = f"{label}_{next_idx}"
                assigned_surya[label] = next_idx + 1
        visuals.append(
            {
                "top": box["top"],
                "left": box["left"],
                "width": box["width"],
                "height": box["height"],
                "label": visual.get("label", "Visual"),
                "visualKey": visual_key,
            }
        )
    density = visual_meta.get("density", "Medium")
    slide["density_proxy"] = density

    # Also include extracted images/tables/charts as visual overlays
    for img_data in slide.get("images", []):
        box = _element_box_to_percent(img_data, slide)
        visuals.append(
            {
                "top": box["top"],
                "left": box["left"],
                "width": box["width"],
                "height": box["height"],
                "label": f"Image: {img_data.get('id', 'img')}",
                "visualKey": img_data.get("id", "img"),
            }
        )
    for tbl_data in slide.get("tables", []):
        if tbl_data.get("x", 0) > 0 or tbl_data.get("y", 0) > 0:
            box = _element_box_to_percent(tbl_data, slide)
            visuals.append(
                {
                    "top": box["top"],
                    "left": box["left"],
                    "width": box["width"],
                    "height": box["height"],
                    "label": f"Table: {tbl_data.get('title', 'tbl')}",
                    "visualKey": tbl_data.get("id", "tbl"),
                }
            )
    for chart_data in slide.get("charts", []):
        if chart_data.get("x", 0) > 0 or chart_data.get("y", 0) > 0:
            box = _element_box_to_percent(chart_data, slide)
            visuals.append(
                {
                    "top": box["top"],
                    "left": box["left"],
                    "width": box["width"],
                    "height": box["height"],
                    "label": f"Chart: {chart_data.get('title', 'chart')}",
                    "visualKey": chart_data.get("id", "chart"),
                }
            )

    # --- Map annotation fixes to REAL bounding boxes ---
    fixes = []
    for a in slide_annotations:
        if a.get("severity") not in ("warning", "hard_block"):
            continue

        # Try to find the source element's bounding box
        fix_box = _annotation_to_bounding_box(a, slide)
        fix_box["label"] = a.get("message", "Issue detected")[:80]
        fix_box["severity"] = a.get("severity", "warning")
        fix_box["category"] = a.get("category", "")
        fix_box["suggestion"] = a.get("suggestion", "")
        fix_box["annotation_id"] = _build_annotation_id(a)
        fixes.append(fix_box)

    # Build council debate from persona results
    council = []
    categories = {
        "structure": "Storyteller",
        "claim_grounding": "Data Auditor",
        "claim_extraction": "Data Auditor",
        "data_accuracy": "Chairman",
        "visual": "Designer",
        "language": "Storyteller",
        "grammar": "Storyteller",
        "hedging": "Chairman",
        "quality": "Storyteller",
        "tone": "Chairman",
        "framework": "Storyteller",
        "so_what": "Chairman",
        "benchmarking": "Data Auditor",
    }

    # Group annotations by persona to avoid duplicates
    persona_comments = {}
    for a in slide_annotations:
        cat = a.get("category", "").lower()
        persona = categories.get(cat, "Chairman")
        if persona not in persona_comments:
            persona_comments[persona] = {
                "persona": persona,
                "text": a.get("message", ""),
                "sentiment": "critical"
                if a.get("severity") == "hard_block"
                else "negative",
                "score": 3 if a.get("severity") == "hard_block" else 6,
            }

    # Add positive entries for categories with no issues
    for persona in ["Chairman", "Storyteller", "Data Auditor", "Designer"]:
        if persona not in persona_comments:
            persona_comments[persona] = {
                "persona": persona,
                "text": "No significant issues detected in this area.",
                "sentiment": "positive",
                "score": 9,
            }
    council = list(persona_comments.values())

    framework_meta = (
        agent_metadata.get("Framework Identifier Agent", {})
        .get("slides_framework", {})
        .get(str(index), {})
    )
    so_what_meta = (
        agent_metadata.get("So What Test Agent", {})
        .get("slides_so_what", {})
        .get(str(index), {})
    )
    benchmark_meta = (
        agent_metadata.get("Competitive Benchmark Agent", {})
        .get("slides_benchmarking", {})
        .get(str(index), {})
    )
    context_meta = (
        agent_metadata.get("Slide Context Synthesizer", {})
        .get("slides_context", {})
        .get(str(index), {})
    )

    framework = framework_meta.get("framework")
    if not framework:
        full_text = slide.get("full_text", "").upper()
        for fallback in [
            "SWOT",
            "PESTEL",
            "PORTER",
            "BCG",
            "VALUE CHAIN",
            "ANSOFF",
            "MECE",
        ]:
            if fallback in full_text:
                framework = fallback
                break

    consultant_score = _build_slide_consultant_score(
        slide,
        slide_annotations,
        deep_analysis_by_slide.get(str(index), {}),
        visual_meta,
    )
    reliability = _build_slide_reliability(
        slide,
        slide_annotations,
        deep_analysis_by_slide.get(str(index), {}),
        visual_meta,
    )
    slide_score = consultant_score["overall_score"]
    guardrail_coverage = _build_guardrail_coverage(
        session=session,
        slide=slide,
        slide_annotations=slide_annotations,
        slide_score=slide_score,
    )
    active_guardrail = _guardrail_from_session(session)
    dynamic_scorecard = _build_dynamic_slide_scorecard(
        active_guardrail,
        consultant_score,
        reliability,
        slide_annotations,
        guardrail_coverage,
        framework_meta,
        so_what_meta,
        benchmark_meta,
    )
    ocr_backend = str(slide.get("ocr_backend") or "native")
    analysis_backends = {
        "surya": bool(visual_agent_meta.get("surya_used")),
        "vision": visual_agent_meta.get("vision_backend", "unknown"),
        "ocr": ocr_backend,
    }

    return {
        "id": f"slide-{index}",
        "title": slide.get("title", "Untitled"),
        "summary": slide.get("full_text", "")[:300]
        + ("..." if len(slide.get("full_text", "")) > 300 else ""),
        "overallScore": slide_score,
        "scoreBreakdown": consultant_score["breakdown"],
        "consultantSummary": consultant_score["consultant_summary"],
        "visualCoverage": consultant_score["visual_coverage"],
        "reliability": reliability,
        "imageAnalysis": keyed_image_analysis,
        "density": density,
        "visuals": visuals,
        "fixes": fixes,
        "councilDebate": council,
        "frameworkDetected": framework,
        "frameworkAnalysis": framework_meta or None,
        "soWhatResult": so_what_meta or None,
        "benchmarkAnalysis": benchmark_meta or None,
        "slideContext": context_meta or None,
        "citationIssues": [
            a.get("message")
            for a in slide_annotations
            if "citation" in a.get("category", "") or "claim" in a.get("category", "")
        ],
        "imageCount": len(slide.get("images", [])),
        "tableCount": len(slide.get("tables", [])),
        "chartCount": len(slide.get("charts", [])),
        "deepAnalysis": deep_analysis_by_slide.get(
            str(index),
            {
                "agents": [],
                "judge": {
                    "name": "Language Analysis",
                    "findings": [],
                    "summary": "No deep analysis available.",
                },
            },
        ),
        "guardrailCoverage": guardrail_coverage,
        "analysisBackends": analysis_backends,
        "dynamicScorecard": dynamic_scorecard,
    }


@app.post("/api/session/{session_id}/slide/{index}/deep-analysis")
async def rerun_slide_deep_analysis(session_id: str, index: int):
    session = _get_session_or_404(session_id)
    session["session_id"] = session_id

    if not session.get("slides_data"):
        raise HTTPException(status_code=400, detail="No slides available for analysis")
    if not session.get("scorecard"):
        raise HTTPException(
            status_code=400,
            detail="Run full analysis before requesting slide-only deep analysis",
        )

    deep_payload = await _refresh_single_slide_deep_analysis(session, index)
    _persist_session(session_id, session)
    return {
        "status": "success",
        "session_id": session_id,
        "slide_index": index,
        "deepAnalysis": deep_payload,
    }


@app.get("/api/session/{session_id}/scorecard")
async def get_scorecard(session_id: str):
    session = _get_session_or_404(session_id)
    scorecard = session.get("scorecard")

    if not scorecard:
        raise HTTPException(status_code=400, detail="No analysis completed")

    annotations = scorecard.get("annotations", [])
    scorecard["annotations"] = [
        {**annotation, "annotation_id": _build_annotation_id(annotation)}
        for annotation in annotations
    ]
    return scorecard


@app.get("/api/session/{session_id}/guardrail")
async def get_session_guardrail(session_id: str):
    ensure_services_loaded()
    session = _get_session_or_404(session_id)
    guardrail = session.get("guardrail")
    if guardrail is None:
        guardrail = guardrail_manager.create_guardrail(
            client_namespace=session.get("client_namespace")
        )
        session["guardrail"] = guardrail
        _persist_session(session_id, session)

    return guardrail.model_dump()


@app.post("/api/session/{session_id}/guardrail/apply")
async def apply_session_guardrail(
    session_id: str,
    guardrail: GuardrailSchema,
    x_user_role: Optional[str] = Header(None),
):
    ensure_services_loaded()
    session = _get_session_or_404(session_id)
    applied_guardrail, signature_cleared = _apply_guardrail_to_session_state(
        session, guardrail
    )
    _invalidate_session_analysis(session)
    _persist_session(session_id, session)

    playbook_rule_count = len(applied_guardrail.playbook_rules or [])
    human_rule_count = len(applied_guardrail.human_confirmed_rules or [])
    audit_log_service.log_event(
        session_id=session_id,
        user_role=x_user_role or "junior",
        action="APPLY_GUARDRAIL",
        details={
            "engagement_type": applied_guardrail.engagement_type,
            "client_namespace": applied_guardrail.client_namespace,
            "playbook_rule_count": playbook_rule_count,
            "human_confirmed_rule_count": human_rule_count,
            "signature_cleared": signature_cleared,
        },
    )

    return {
        "status": "applied",
        "guardrail": applied_guardrail.model_dump(),
        "analysis_invalidated": True,
    }


@app.post("/api/session/{session_id}/guardrail/template/save")
async def save_session_guardrail_template(
    session_id: str,
    req: SaveTemplateRequest,
    x_user_role: Optional[str] = Header(None),
):
    ensure_services_loaded()
    session = _get_session_or_404(session_id)
    guardrail = session.get("guardrail")
    if not isinstance(guardrail, GuardrailSchema):
        raise HTTPException(status_code=400, detail="No guardrail in session")

    template_guardrail = GuardrailSchema(**guardrail.model_dump())
    template_guardrail.signed_by = None
    template_guardrail.signed_at = None
    template_guardrail.sha256 = None
    template_guardrail.signature = None
    template_guardrail.public_key = None
    template_guardrail.signature_algorithm = None

    slug = _slugify_guardrail_template_name(req.template_name)
    namespace = (
        template_guardrail.client_namespace
        or session.get("client_namespace")
        or "shared"
    )
    engagement = template_guardrail.engagement_type or "strategy"
    timestamp = utc_now_compact()
    filename = f"template_{slug}_{namespace}_{engagement}_{timestamp}.json"
    filepath = _guardrail_template_dir() / filename
    with open(filepath, "w") as f:
        json.dump(template_guardrail.model_dump(), f, indent=2)

    audit_log_service.log_event(
        session_id=session_id,
        user_role=x_user_role or "junior",
        action="SAVE_GUARDRAIL_TEMPLATE",
        details={
            "template_name": req.template_name.strip(),
            "filename": filename,
            "engagement_type": engagement,
            "client_namespace": namespace,
        },
    )

    return {
        "status": "saved",
        "template": {
            "id": filename,
            "filename": filename,
            "template_name": req.template_name.strip(),
            "engagement_type": engagement,
            "client_namespace": namespace,
        },
    }


@app.get("/api/guardrail/template/list")
async def list_guardrail_templates(session_id: Optional[str] = None):
    ensure_services_loaded()
    template_dir = _guardrail_template_dir()
    files = sorted(
        template_dir.glob("*.json"), key=lambda f: f.stat().st_mtime, reverse=True
    )

    session_guardrail = None
    client_namespace = None
    engagement_type = None
    if session_id:
        session = _get_session_or_404(session_id)
        session_guardrail = session.get("guardrail")
        client_namespace = session.get("client_namespace")

    templates = []
    for path in files:
        try:
            template_guardrail = guardrail_manager.load_guardrail(
                str(path), require_signature=False
            )
        except Exception:
            continue

        template_name = path.stem
        if template_name.startswith("template_"):
            template_name = template_name[len("template_") :]

        guardrail_namespace = template_guardrail.client_namespace or "shared"
        guardrail_engagement = template_guardrail.engagement_type or "strategy"
        current_engagement = (
            session_guardrail.engagement_type
            if isinstance(session_guardrail, GuardrailSchema)
            else engagement_type
        )
        compatible = True
        if (
            current_engagement
            and guardrail_engagement
            and current_engagement != guardrail_engagement
        ):
            compatible = False

        templates.append(
            {
                "id": path.name,
                "filename": path.name,
                "template_name": template_name.replace("_", " "),
                "engagement_type": guardrail_engagement,
                "client_namespace": template_guardrail.client_namespace,
                "rule_count": len(template_guardrail.playbook_rules or [])
                + len(template_guardrail.human_confirmed_rules or []),
                "compatible": compatible,
                "updated_at": datetime.utcfromtimestamp(
                    path.stat().st_mtime
                ).isoformat()
                + "Z",
                "scope": (
                    "shared"
                    if guardrail_namespace == "shared"
                    else (
                        "namespace"
                        if client_namespace and client_namespace == guardrail_namespace
                        else "custom"
                    )
                ),
            }
        )

    return {"templates": templates}


@app.post("/api/session/{session_id}/guardrail/template/{template_id}/activate")
async def activate_guardrail_template(
    session_id: str,
    template_id: str,
    x_user_role: Optional[str] = Header(None),
):
    ensure_services_loaded()
    session = _get_session_or_404(session_id)
    template_guardrail = guardrail_manager.load_guardrail(
        str(_resolve_guardrail_template_file(template_id)),
        require_signature=False,
    )
    applied_guardrail, signature_cleared = _apply_guardrail_to_session_state(
        session, template_guardrail
    )
    _invalidate_session_analysis(session)
    _persist_session(session_id, session)

    audit_log_service.log_event(
        session_id=session_id,
        user_role=x_user_role or "junior",
        action="ACTIVATE_GUARDRAIL_TEMPLATE",
        details={
            "template_id": template_id,
            "engagement_type": applied_guardrail.engagement_type,
            "client_namespace": applied_guardrail.client_namespace,
            "signature_cleared": signature_cleared,
        },
    )

    return {
        "status": "activated",
        "guardrail": applied_guardrail.model_dump(),
        "analysis_invalidated": True,
    }


@app.get("/api/session/{session_id}/slides")
async def get_slides(session_id: str):
    session = _get_session_or_404(session_id)
    await _hydrate_slide_assets_for_session(session_id, session)
    slides_data = session.get("slides_data", [])

    return {"slides": slides_data}


@app.get("/api/session/{session_id}/download")
async def download_annotated(session_id: str):
    session = _get_session_or_404(session_id)
    deck_path = session.get("deck_path")

    if not deck_path or not os.path.exists(deck_path):
        raise HTTPException(status_code=400, detail="No deck uploaded")

    scorecard_data = session.get("scorecard")
    if not scorecard_data:
        raise HTTPException(status_code=400, detail="No analysis completed")

    ensure_services_loaded()
    annotations = [Annotation(**a) for a in scorecard_data.get("annotations", [])]

    try:
        output_path = language_agent.write_comments_to_pptx(deck_path, annotations)

        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="SlideForge_annotated.pptx",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/guardrail/create")
async def create_guardrail(
    engagement_type: str = "strategy",
    client_namespace: Optional[str] = None,
    playbook_rules: list[dict] = [],
    human_confirmed_rules: list[dict] = [],
):
    ensure_services_loaded()
    guardrail = guardrail_manager.create_guardrail(
        engagement_type=engagement_type,
        client_namespace=client_namespace,
        playbook_rules=playbook_rules,
        human_confirmed_rules=human_confirmed_rules,
    )

    return guardrail.model_dump()


@app.post("/api/guardrail/{session_id}/sign")
async def sign_guardrail(
    session_id: str, req: SignRequest, role: str = Depends(require_senior)
):
    ensure_services_loaded()
    session = _get_session_or_404(session_id)
    guardrail = session.get("guardrail")

    if not guardrail:
        raise HTTPException(status_code=400, detail="No guardrail in session")

    signed_guardrail = guardrail_manager.sign_guardrail(guardrail, req.user_name)

    filepath = guardrail_manager.save_guardrail(signed_guardrail)

    session["guardrail"] = signed_guardrail
    _persist_session(session_id, session)

    audit_log_service.log_event(
        session_id=session_id,
        user_role="senior",
        action="SIGN_GUARDRAIL",
        details={"signer": req.user_name, "filepath": filepath},
    )

    return {
        "signed_guardrail": signed_guardrail.model_dump(),
        "filepath": filepath,
    }


@app.get("/api/session/{session_id}/audit-log")
async def get_audit_log(session_id: str):
    _get_session_or_404(session_id)
    try:
        events = audit_log_service.get_events(session_id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to load audit log: {e}")
    return {"entries": events}


@app.get("/api/guardrail/list")
async def list_guardrails(session_id: str):
    guardrail_dir = data_dir / "guardrails"

    if not guardrail_dir.exists():
        return {"guardrails": []}

    session = _get_session_or_404(session_id)
    guardrail = session.get("guardrail")
    namespace = (
        (guardrail.client_namespace if guardrail else None)
        or session.get("client_namespace")
        or "default"
    )
    engagement = (guardrail.engagement_type if guardrail else None) or "strategy"
    prefix = f"guardrail_{namespace}_{engagement}_"
    files = sorted(f for f in guardrail_dir.glob("*.json") if f.name.startswith(prefix))

    return {"guardrails": [{"filename": f.name, "id": f.name} for f in files]}


async def get_current_role(x_user_role: Optional[str] = Header(None)):
    return x_user_role or "junior"


@app.post("/api/session/{session_id}/override")
async def record_override(
    session_id: str, req: OverrideRequest, role: str = Depends(get_current_role)
):
    session = _get_session_or_404(session_id)
    scorecard_data = session.get("scorecard", {})

    # Resolve annotation details from scorecard
    # We use 'message' as the ID coming from frontend (it's what we have)
    annotation = {}
    if scorecard_data:
        for a in scorecard_data.get("annotations", []):
            candidate_ids = {
                _build_annotation_id(a),
                a.get("message", ""),
                a.get("text", ""),
            }
            if req.annotation_id in candidate_ids:
                annotation = a
                break

    category = annotation.get("category", "unknown")
    slide_index = annotation.get("slide_index", 0)

    adaptation_agent.log_override_decision(
        session_id=session_id,
        annotation_id=req.annotation_id,
        category=category,
        reason=req.reason,
        slide_index=slide_index,
    )

    audit_log_service.log_event(
        session_id=session_id,
        user_role=role,
        action="OVERRIDE",
        details={"id": req.annotation_id, "reason": req.reason, "slide": slide_index},
    )

    if scorecard_data:
        annotations = scorecard_data.get("annotations", [])
        filtered_annotations = []
        for annotation in annotations:
            candidate_ids = {
                _build_annotation_id(annotation),
                annotation.get("message", ""),
                annotation.get("text", ""),
            }
            if req.annotation_id not in candidate_ids:
                filtered_annotations.append(annotation)
        annotations = filtered_annotations
        scorecard_data["annotations"] = annotations
        session["scorecard"] = scorecard_data
        _refresh_scorecard_counters(session)
        session["annotations_by_slide"] = _group_annotations_by_slide(annotations)
        _remove_annotation_from_deep_analysis(session, req.annotation_id)
        session["status"] = "analyzed"
        _persist_session(session_id, session)

    return {"status": "recorded", "session_id": session_id}


@app.post("/api/session/{session_id}/accept")
async def accept_fix(
    session_id: str, req: AcceptRequest, role: str = Depends(get_current_role)
):
    session = _get_session_or_404(session_id)
    scorecard_data = session.get("scorecard")

    if not scorecard_data:
        raise HTTPException(status_code=400, detail="No scorecard found")

    annotations = scorecard_data.get("annotations", [])
    accepted = []
    applied = []
    for annotation in annotations:
        candidate_ids = {
            _build_annotation_id(annotation),
            annotation.get("message", ""),
            annotation.get("text", ""),
        }
        if req.annotation_id in candidate_ids:
            applied.append(annotation)
        else:
            accepted.append(annotation)

    audit_log_service.log_event(
        session_id=session_id,
        user_role=role,
        action="ACCEPT_FIX",
        details={"id": req.annotation_id},
    )

    scorecard_data["annotations"] = accepted
    session["scorecard"] = scorecard_data
    _refresh_scorecard_counters(session)
    session["annotations_by_slide"] = _group_annotations_by_slide(accepted)
    _remove_annotation_from_deep_analysis(session, req.annotation_id)
    session["status"] = "analyzed"
    _persist_session(session_id, session)

    return {"status": "accepted", "applied_count": len(applied)}


@app.post("/api/session/{session_id}/revision")
async def run_revision_loop(session_id: str):
    session = _get_session_or_404(session_id)
    scorecard_data = session.get("scorecard")
    slides_data = session.get("slides_data", [])

    if not scorecard_data:
        raise HTTPException(status_code=400, detail="No analysis completed")

    scorecard = QAScorecard(**scorecard_data)

    attempt = 0
    max_attempts = 3
    score_history = [scorecard.composite_score]

    while attempt < max_attempts:
        should_revise = await revision_orchestrator.should_revise(scorecard, attempt)

        if not should_revise:
            break

        fixes = await revision_orchestrator.apply_auto_remediation(
            scorecard, slides_data
        )

        session["auto_fixes"] = session.get("auto_fixes", []) + fixes

        current_guardrail = session.get("guardrail", GuardrailSchema())
        if isinstance(current_guardrail, dict):
            current_guardrail = GuardrailSchema(**current_guardrail)
            session["guardrail"] = current_guardrail
        excel_data = session.get("excel_data")

        agent_results = await analysis_orchestrator.run_parallel_analysis(
            slides_data, current_guardrail, excel_data
        )

        language_annotations = await language_agent.analyze_deck(
            slides_data, current_guardrail.language_rules
        )
        source_grounding_annotations = await _build_source_grounding_annotations(
            session
        )
        language_annotations.extend(source_grounding_annotations)

        context_result = await analysis_orchestrator.run_slide_context_synthesizer(
            slides_data,
            agent_results,
            language_annotations,
        )
        agent_results.append(context_result)

        scorecard = qa_grader.calculate_scorecard(
            agent_results,
            language_annotations,
            current_guardrail,
        )

        score_history.append(scorecard.composite_score)
        scorecard_data = scorecard.model_dump()
        scorecard_data["weights"] = _normalize_guardrail_weights(
            session.get("guardrail", GuardrailSchema())
        )
        session["scorecard"] = scorecard_data
        session["annotations_by_slide"] = _group_annotations_by_slide(
            scorecard_data.get("annotations", [])
        )
        deep_analysis_by_slide = _build_deep_analysis_by_slide(
            agent_results, language_annotations, slides_data
        )
        deep_analysis_by_slide = await _enrich_deep_analysis_with_slide_reviews(
            slides_data, deep_analysis_by_slide, current_guardrail
        )
        session["deep_analysis_by_slide"] = deep_analysis_by_slide

        agent_metadata = {}
        for res in agent_results:
            agent_metadata[res.agent_name] = res.metadata
        session["agent_metadata"] = agent_metadata
        session["status"] = "analyzed"
        _persist_session(session_id, session)

        attempt += 1

    return {
        "session_id": session_id,
        "revision_count": attempt,
        "score_history": score_history,
        "final_score": scorecard.composite_score,
        "auto_fixes_applied": len(session.get("auto_fixes", [])),
    }


@app.post("/api/session/{session_id}/prepare")
async def prepare_for_delivery(session_id: str, role: str = Depends(require_senior)):
    session = _get_session_or_404(session_id)

    # GAP-12: Enforce Sign-off
    if not session.get("senior_signed"):
        raise HTTPException(
            status_code=400,
            detail="Senior sign-off mandatory before delivery packaging",
        )

    deck_path = session.get("deck_path")
    if not deck_path or not os.path.exists(deck_path):
        raise HTTPException(status_code=400, detail="No deck uploaded")

    from pptx import Presentation

    if deck_path.endswith(".pptx"):
        prs = Presentation(deck_path)
        for slide in prs.slides:
            for comment in list(slide.comments):
                if "SlideForge" in comment.author:
                    slide.comments._comments.remove(comment._element)

        prs.core_properties.title = ""
        prs.core_properties.subject = ""
        prs.core_properties.keywords = ""
        prs.core_properties.comments = ""

        # Scrub hidden notes
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_notes_slide:
                    notes = shape.notes_slide.notes_text_frame
                    if "SlideForge" in notes.text:
                        notes.text = ""

        output_path = deck_path.replace(".pptx", "_delivery.pptx")
        prs.save(output_path)
    else:
        output_path = deck_path.replace(".pdf", "_delivery.pdf")
        shutil.copy(deck_path, output_path)

    # GAP-07/09/12: Create Delivery Package ZIP
    package_path = Path(deck_path).parent / "Delivery_Package.zip"
    audit_csv_path = Path(deck_path).parent / "Audit_Log.csv"
    metadata_path = Path(deck_path).parent / "Metadata.json"

    # Export Audit Log
    audit_log_service.export_csv(session_id, str(audit_csv_path))

    # Create Metadata
    metadata = {
        "session_id": session_id,
        "delivery_timestamp": utc_now_iso(),
        "compliance_status": "PASSED",
        "senior_signatory": session.get("senior_name", "N/A"),
        "composite_score": session.get("scorecard", {}).get("composite_score", 0),
        "hard_blocks_remaining": 0,
        "artifacts": [os.path.basename(output_path), "Audit_Log.csv"],
    }
    with open(metadata_path, "w") as f:
        json.dump(metadata, f, indent=4)

    with zipfile.ZipFile(package_path, "w") as zipf:
        zipf.write(output_path, arcname=os.path.basename(output_path))
        zipf.write(audit_csv_path, arcname="Audit_Log.csv")
        zipf.write(metadata_path, arcname="Metadata.json")

    return {
        "status": "ready",
        "download_url": f"/api/session/{session_id}/download-package",
        "metadata": metadata,
    }


@app.post("/api/session/{session_id}/sign-off")
async def sign_off_session(
    session_id: str, request: SignRequest, role: str = Depends(require_senior)
):
    session = _get_session_or_404(session_id)
    session["senior_signed"] = True
    session["senior_name"] = request.user_name

    audit_log_service.log_event(
        session_id=session_id,
        user_role="senior",
        action="SENIOR_SIGN_OFF",
        details={"signer": request.user_name},
    )

    _persist_session(session_id, session)

    return {"status": "signed", "signer": request.user_name}


@app.get("/api/session/{session_id}/delivery-status")
async def get_delivery_status(session_id: str):
    session = _get_session_or_404(session_id)
    return {
        "senior_signed": bool(session.get("senior_signed")),
        "senior_name": session.get("senior_name"),
    }


@app.get("/api/session/{session_id}/download-package")
async def download_package(session_id: str):
    session = _get_session_or_404(session_id)
    deck_path = session.get("deck_path")
    package_path = Path(deck_path).parent / "Delivery_Package.zip"
    if not package_path.exists():
        raise HTTPException(status_code=404, detail="Package not ready")
    return FileResponse(
        package_path, filename="Delivery_Package.zip", media_type="application/zip"
    )


@app.get("/api/guardrail/diff")
async def diff_guardrails(session_id: str, old_id: str, new_id: str):
    ensure_services_loaded()
    _get_session_or_404(session_id)
    guardrail_dir = data_dir / "guardrails"

    def _resolve_guardrail_file(file_id: str) -> Path:
        candidate = (guardrail_dir / Path(file_id).name).resolve()
        guardrail_root = guardrail_dir.resolve()
        if guardrail_root not in candidate.parents or not candidate.exists():
            raise HTTPException(status_code=400, detail="Invalid guardrail reference")
        return candidate

    try:
        old_guardrail = guardrail_manager.load_guardrail(
            str(_resolve_guardrail_file(old_id))
        )
        new_guardrail = guardrail_manager.load_guardrail(
            str(_resolve_guardrail_file(new_id))
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

    diff = guardrail_manager.diff_guardrails(old_guardrail, new_guardrail)

    return {
        "old_version": old_guardrail.schema_version,
        "new_version": new_guardrail.schema_version,
        "diff": diff,
    }


@app.post("/api/patterns/log-engagement")
async def log_engagement_completion(
    session_id: str,
    engagement_type: str,
    client_namespace: str = None,
):
    session = _get_session_or_404(session_id)
    scorecard_data = session.get("scorecard", {})

    adaptation_agent.log_engagement_completion(
        engagement_type=engagement_type,
        client_namespace=client_namespace,
        score_history=scorecard_data.get("score_history", []),
        resolved_issues=[],
        persistent_issues=[],
        revision_count=scorecard_data.get("revision_count", 0),
        final_score=scorecard_data.get("composite_score", 0),
    )

    return {"status": "logged"}


@app.get("/api/patterns/suggestions")
async def get_pattern_suggestions(engagement_type: str = None):
    analysis = adaptation_agent.analyze_and_suggest(engagement_type)

    return analysis


@app.post("/api/template/discover/start")
async def start_discovery(session_id: str, slides_data: Optional[list[dict]] = None):
    if session_id not in active_sessions:
        # Create temp session if not exists
        now_ts = time.time()
        namespace_key = _namespace_key(None)
        _enforce_session_capacity(namespace_key)
        active_sessions[session_id] = {
            "id": session_id,
            "client_namespace": None,
            "namespace_key": namespace_key,
            "status": "discovery",
            "created_at_ts": now_ts,
            "last_access_ts": now_ts,
        }
        _persist_session(session_id, active_sessions[session_id])

    session = _get_session_or_404(session_id)
    from app.agents.template_discovery import template_discovery_agent

    # 1. Initial observation
    gold_slides = slides_data or session.get("slides_data", [])
    if not gold_slides:
        raise HTTPException(status_code=400, detail="No slides available for discovery")
    patterns = await template_discovery_agent._observe_slides(gold_slides)

    # 2. Initialize discovery state
    session["discovery_state"] = {
        "patterns": patterns,
        "discovered_patterns": patterns,
        "questions": [],
        "answers": [],
        "current_question": None,
        "step": "questioning",
        "gold_slides": gold_slides,
    }
    _persist_session(session_id, session)

    # 3. Get first question
    question = await template_discovery_agent._generate_next_question(
        {"discovered_patterns": patterns, "questions": [], "answers": []}, []
    )

    if question:
        session["discovery_state"]["current_question"] = {
            "question": question.question,
            "evidence": question.evidence,
            "context": question.context,
        }
        _persist_session(session_id, session)
        return {
            "status": "questioning",
            "question": session["discovery_state"]["current_question"],
        }
    else:
        # Jump to complete if no questions
        schema = template_discovery_agent._build_schema(session["discovery_state"])
        session["discovery_state"]["step"] = "completed"
        session["discovery_state"]["schema"] = schema
        _persist_session(session_id, session)
        return {"status": "completed", "schema": schema}


@app.post("/api/template/discover/answer")
async def answer_discovery(session_id: str, answer: str):
    session = _get_session_or_404(session_id)
    state = session.get("discovery_state")
    if not state:
        raise HTTPException(status_code=400, detail="Discovery not started")
    if not state.get("current_question"):
        raise HTTPException(status_code=400, detail="No active discovery question")

    # Record answer
    current_q = state["current_question"]
    state["answers"].append(
        {
            "question": current_q["question"],
            "answer": answer,
            "context": current_q["context"],
        }
    )
    state["questions"].append(current_q)
    _persist_session(session_id, session)

    from app.agents.template_discovery import template_discovery_agent

    # Limit to 3 questions for this POC/GAP-08 requirement
    if len(state["answers"]) >= 3:
        schema = template_discovery_agent._build_schema(state)
        state["step"] = "completed"
        state["schema"] = schema
        _persist_session(session_id, session)
        return {"status": "completed", "schema": schema}

    # Get next question
    next_q = await template_discovery_agent._generate_next_question(
        {
            "discovered_patterns": state.get("discovered_patterns")
            or state["patterns"],
            "questions": state["questions"],
            "answers": state["answers"],
        },
        [],
    )

    if next_q:
        state["current_question"] = {
            "question": next_q.question,
            "evidence": next_q.evidence,
            "context": next_q.context,
        }
        _persist_session(session_id, session)
        return {"status": "questioning", "question": state["current_question"]}
    else:
        schema = template_discovery_agent._build_schema(state)
        state["step"] = "completed"
        state["schema"] = schema
        _persist_session(session_id, session)
        return {"status": "completed", "schema": schema}


@app.post("/api/template/discover")
async def discover_template(
    gold_slides: list[dict] = None,
    playbook_text: str = None,
    engagement_type: str = "strategy",
):
    if not gold_slides and not playbook_text:
        raise HTTPException(
            status_code=400, detail="Provide gold_slides or playbook_text"
        )

    if gold_slides and playbook_text:
        guardrail = await template_discovery_agent.discover_combined(
            gold_slides, playbook_text
        )
    elif gold_slides:
        guardrail = await template_discovery_agent.discover_from_gold_slides(
            gold_slides
        )
    else:
        guardrail = await template_discovery_agent.discover_from_playbook(playbook_text)

    guardrail.engagement_type = engagement_type

    return guardrail.model_dump()


@app.post("/api/template/discover/upload")
async def discover_template_from_upload(
    session_id: str, file: UploadFile = File(...), engagement_type: str = "strategy"
):
    _get_session_or_404(session_id)

    allowed_exts = {".pdf", ".docx", ".txt", ".md"}
    filename = Path(file.filename or "playbook.txt").name
    ext = os.path.splitext(filename)[1].lower()
    if ext not in allowed_exts:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported playbook extension: {ext}. Allowed: {sorted(allowed_exts)}",
        )

    content = await file.read()
    _enforce_upload_size(content)
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded playbook file is empty")

    try:
        extracted_documents = _extract_uploaded_document_text(filename, content)
    except Exception as exc:
        raise HTTPException(
            status_code=500, detail=f"Failed to extract playbook text: {exc}"
        ) from exc

    playbook_text = "\n\n".join(
        doc for doc in extracted_documents if doc.strip()
    ).strip()
    if not playbook_text:
        raise HTTPException(
            status_code=400,
            detail="No readable text could be extracted from the uploaded playbook",
        )

    guardrail = await template_discovery_agent.discover_from_playbook(playbook_text)
    guardrail.engagement_type = engagement_type

    audit_log_service.log_event(
        session_id=session_id,
        user_role="analyst",
        action="DISCOVERY_PLAYBOOK_UPLOAD",
        details={
            "filename": filename,
            "engagement_type": engagement_type,
            "extracted_characters": len(playbook_text),
        },
    )

    return {
        "status": "completed",
        "filename": filename,
        "extracted_characters": len(playbook_text),
        "schema": guardrail.model_dump(),
    }


async def get_local_llm_status():
    response = _runtime_provider_config_response()
    response["analysis_max_tokens"] = inference_service.analysis_max_tokens
    return response


async def update_local_llm_config(payload: RuntimeProviderConfigPayload):
    try:
        new_provider = InferenceProvider(payload.provider)
    except ValueError:
        raise HTTPException(status_code=400, detail="Unsupported provider")

    provider_updates = _extract_provider_updates(payload)

    if new_provider in (
        InferenceProvider.API,
        InferenceProvider.OLLAMA,
        InferenceProvider.LM_STUDIO,
    ):
        current_config = inference_service.get_provider_connection_config(new_provider)
        candidate_config = {
            "base_url": current_config.base_url,
            "api_key": current_config.api_key,
            "model": current_config.model,
            **(provider_updates or {}),
        }

        if not candidate_config.get("base_url"):
            raise HTTPException(
                status_code=400,
                detail="Provider base URL is required.",
            )

        if not candidate_config.get("model"):
            raise HTTPException(
                status_code=400,
                detail="Model is required.",
            )

        if new_provider == InferenceProvider.API and not candidate_config.get(
            "api_key"
        ):
            raise HTTPException(
                status_code=400,
                detail="Cloud AI requires an API key. Add one before applying this provider.",
            )
    elif provider_updates:
        raise HTTPException(
            status_code=400,
            detail="This provider does not support runtime web configuration.",
        )

    update_kwargs = {
        "provider": new_provider,
        "local_context_window": payload.local_context_window,
    }
    if new_provider == InferenceProvider.API:
        update_kwargs["api"] = provider_updates
    elif new_provider == InferenceProvider.OLLAMA:
        update_kwargs["ollama"] = provider_updates
    elif new_provider == InferenceProvider.LM_STUDIO:
        update_kwargs["lm_studio"] = provider_updates

    try:
        inference_service.update_runtime_provider_config(**update_kwargs)
        inference_service.reinitialize_provider()
        response = _runtime_provider_config_response()
        response["status"] = "success"
        return response
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


async def set_session_llm_settings(
    session_id: str, provider: str, context_window: Optional[int] = None
):
    session = _get_session_or_404(session_id)
    if provider not in (
        InferenceProvider.API.value,
        InferenceProvider.MLX.value,
        InferenceProvider.TRANSFORMERS.value,
        InferenceProvider.OLLAMA.value,
        InferenceProvider.LM_STUDIO.value,
    ):
        raise HTTPException(status_code=400, detail="Unsupported provider")

    normalized_context = None
    if provider in (InferenceProvider.OLLAMA.value, InferenceProvider.LM_STUDIO.value):
        normalized_context = max(
            1024,
            min(int(context_window or inference_service.local_context_window), 131072),
        )

    session["llm_settings"] = {
        "provider": provider,
        "context_window": normalized_context,
    }
    _persist_session(session_id, session)
    return {"status": "success", "llm_settings": session["llm_settings"]}


async def get_analysis_settings():
    return {"analysis_max_tokens": inference_service.analysis_max_tokens}


async def update_analysis_settings(analysis_max_tokens: int):
    try:
        inference_service.set_analysis_max_tokens(analysis_max_tokens)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    return {
        "status": "success",
        "analysis_max_tokens": inference_service.analysis_max_tokens,
    }


async def get_grammar_status():
    lt_status = await language_agent.lt_client.status()
    language_tool_available = lt_status.get("available", False)
    lt_error = lt_status.get("last_error")
    return {
        "enabled": True,
        "engine": "languagetool" if language_tool_available else "regex_fallback",
        "language_tool_available": language_tool_available,
        "base_url": language_agent.lt_client.base_url,
        "check_url": lt_status.get("check_url"),
        "last_error": lt_error,
        "notes": (
            "Full grammar and spelling checks are active."
            if language_tool_available
            else (
                "LanguageTool is unavailable, so lighter regex-based grammar checks are active."
                + (f" Last error: {lt_error}" if lt_error else "")
            )
        ),
    }


async def test_llm_connection():
    if inference_service.llm is None:
        if inference_service.current_provider == InferenceProvider.API.value:
            api_config = inference_service.get_provider_connection_config(
                InferenceProvider.API
            )
            if not api_config.api_key:
                raise HTTPException(
                    status_code=400,
                    detail="Cloud AI is selected but no API key is configured. Save provider settings and try again.",
                )
        raise HTTPException(
            status_code=400,
            detail="No LLM provider configured. Save provider settings and try again.",
        )
    try:
        messages = [Message(role="user", content="Say exactly 'OK' and nothing else.")]
        response = await inference_service.llm.generate(messages, max_tokens=10)
        return {
            "status": "success",
            "response": response.content,
            "provider": inference_service.current_provider,
            "model": response.model,
        }
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


async def get_llm_diagnostics():
    """Connection diagnostics and provider status for the LLM engine."""
    import socket

    diagnostics = {
        "current_provider": inference_service.current_provider,
        "llm_available": inference_service.llm is not None,
        "providers": {},
        "resilience": {
            "failure_window_seconds": getattr(
                inference_service, "_circuit_failure_window_seconds", 60
            ),
            "failure_threshold": getattr(
                inference_service, "_circuit_failure_threshold", 5
            ),
            "open_seconds": getattr(inference_service, "_circuit_open_seconds", 30),
            "provider_stats": (
                inference_service.get_provider_failure_stats()
                if hasattr(inference_service, "get_provider_failure_stats")
                else {}
            ),
        },
    }

    # Check LM Studio
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.settimeout(1)
            lm_studio_up = s.connect_ex(("localhost", 1234)) == 0
    except Exception:
        lm_studio_up = False
    diagnostics["providers"]["lm_studio"] = {
        "port": 1234,
        "reachable": lm_studio_up,
        **_provider_response(InferenceProvider.LM_STUDIO),
    }

    # Check Ollama
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.settimeout(1)
            ollama_up = s.connect_ex(("localhost", 11434)) == 0
    except Exception:
        ollama_up = False
    diagnostics["providers"]["ollama"] = {
        "port": 11434,
        "reachable": ollama_up,
        **_provider_response(InferenceProvider.OLLAMA),
    }

    diagnostics["providers"]["api"] = _provider_response(InferenceProvider.API)

    # Test current LLM if available
    if inference_service.llm is not None:
        try:
            messages = [Message(role="user", content="Reply with exactly: OK")]
            resp = await inference_service.llm.generate(messages, max_tokens=5)
            diagnostics["last_test"] = {
                "status": "success",
                "response_preview": resp.content[:50] if resp.content else "",
                "model": resp.model,
            }
        except Exception as e:
            diagnostics["last_test"] = {
                "status": "error",
                "error": str(e)[:200],
            }

    return diagnostics


async def list_available_models():
    """List models available from the current LLM provider."""
    import httpx

    provider_config = inference_service.get_provider_connection_config()
    base_url = provider_config.base_url or "http://localhost:11434/v1"
    api_key = provider_config.api_key

    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            response = await client.get(
                f"{base_url}/models",
                headers={"Authorization": f"Bearer {api_key}"},
            )
            if response.status_code == 200:
                data = response.json()
                models = data.get("data", [])
                return {
                    "status": "success",
                    "models": [
                        {"id": m.get("id", ""), "name": m.get("id", "")} for m in models
                    ],
                }
            else:
                return {
                    "status": "error",
                    "models": [],
                    "detail": f"HTTP {response.status_code}",
                }
    except Exception as e:
        return {"status": "error", "models": [], "detail": str(e)[:200]}


@app.post("/api/session/{session_id}/upload-source")
async def upload_source_document(session_id: str, file: UploadFile = File(...)):
    """Upload source documents (PDF/XLSX/DOCX) for claim evidence checking."""
    session = _get_session_or_404(session_id)
    ensure_services_loaded()

    # 1. Source MIME Validation & Extension Hardening
    allowed_source_exts = {".pdf", ".xlsx", ".xls", ".docx", ".txt", ".csv", ".md"}
    filename = Path(file.filename or "source.pdf").name
    ext = os.path.splitext(filename)[1].lower()

    if ext not in allowed_source_exts:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported source extension: {ext}. Allowed: {allowed_source_exts}",
        )

    # 2. Extract Source Content

    content = await file.read()
    _enforce_upload_size(content)
    source_dir = data_dir / "sessions" / session_id / "sources"
    source_dir.mkdir(parents=True, exist_ok=True)

    filepath = source_dir / filename
    filepath.write_bytes(content)

    # Extract text based on file type and index in ChromaDB
    try:
        documents = _extract_uploaded_document_text(filename, content, filepath)
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Failed to parse source document: {e}"
        )

    if documents:
        namespace = f"session_{session_id}"
        chroma_manager.add_documents(
            namespace=namespace,
            documents=documents,
            ids=[f"src_{filename}_{i}" for i in range(len(documents))],
            metadatas=[{"source": filename, "index": i} for i in range(len(documents))],
        )
        session["source_namespace"] = namespace
        session.setdefault("source_files", []).append(filename)
        source_indexed_chunks = session.setdefault("source_indexed_chunks", {})
        source_indexed_chunks[filename] = len(documents)

    _persist_session(session_id, session)

    audit_log_service.log_event(
        session_id=session_id,
        user_role="analyst",
        action="UPLOAD_SOURCE",
        details={"filename": filename, "docs_indexed": len(documents)},
    )

    return {
        "status": "indexed",
        "filename": filename,
        "documents_indexed": len(documents),
        "total_sources": len(session.get("source_files", [])),
    }


@app.post("/api/session/{session_id}/upload-excel")
async def upload_excel_for_data_lineage(session_id: str, file: UploadFile = File(...)):
    """Upload Excel file for data lineage verification (chart vs Excel comparison)."""
    session = _get_session_or_404(session_id)

    filename = Path(file.filename or "data.xlsx").name
    ext = os.path.splitext(filename)[1].lower()
    if ext not in (".xlsx", ".xls"):
        raise HTTPException(
            status_code=400,
            detail="Only Excel files (.xlsx, .xls) are accepted for data lineage.",
        )

    content = await file.read()
    _enforce_upload_size(content)
    upload_dir = data_dir / "uploads" / session_id
    upload_dir.mkdir(parents=True, exist_ok=True)
    filepath = upload_dir / filename
    filepath.write_bytes(content)

    from openpyxl import load_workbook

    try:
        wb = load_workbook(str(filepath), data_only=True)
        sheets = {}
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):
                    rows.append(row_data)
            sheets[sheet.title] = rows

        excel_data = {"sheets": sheets, "file_path": str(filepath)}
        session["excel_data"] = excel_data

        _persist_session(session_id, session)

        audit_log_service.log_event(
            session_id=session_id,
            user_role="analyst",
            action="UPLOAD_EXCEL",
            details={
                "filename": filename,
                "sheet_count": len(sheets),
                "total_rows": sum(len(r) for r in sheets.values()),
            },
        )

        return {
            "status": "loaded",
            "filename": filename,
            "sheets": list(sheets.keys()),
            "total_rows": sum(len(r) for r in sheets.values()),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse Excel: {e}")


@app.get("/api/session/{session_id}/slide/{index}/image-analysis")
async def get_image_analysis(session_id: str, index: int):
    """Return vision model analysis for images in a specific slide."""
    session = _get_session_or_404(session_id)
    agent_metadata = session.get("agent_metadata", {})
    visual_meta = (
        agent_metadata.get("Visual Analysis Agent", {})
        .get("slides_analysis", {})
        .get(str(index), {})
    )

    return {
        "slide_index": index,
        "image_analysis": visual_meta.get("image_analysis", []),
        "visuals": visual_meta.get("visuals", []),
        "density": visual_meta.get("density", "Medium"),
    }


@app.get("/{full_path:path}")
async def serve_spa(full_path: str):
    if full_path.startswith("api/"):
        raise HTTPException(status_code=404, detail="API route not found")

    index_file = static_dir / "index.html"
    if index_file.exists():
        return FileResponse(index_file)
    return {"message": "SlideForge AI Engine", "status": "Waiting for frontend build"}


async def shutdown_event():
    for sid, session in list(active_sessions.items()):
        _persist_session(sid, session)
    try:
        await inference_service.close()
    except Exception:
        logger.exception("Failed to close LLM HTTP client")


if __name__ == "__main__":
    import uvicorn

    parser = argparse.ArgumentParser(description="SlideForge backend server")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8002)
    parser.add_argument("--data-dir", default=None)
    args = parser.parse_args()

    if args.data_dir:
        data_dir = Path(args.data_dir).expanduser().resolve()
        data_dir.mkdir(parents=True, exist_ok=True)
        settings.data_dir = str(data_dir)
        session_store = SQLiteSessionStore(str(data_dir / "sessions.db"))

    uvicorn.run(app, host=args.host, port=args.port)
